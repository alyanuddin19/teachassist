"""
File parser utility - extracts text content AND embedded images from
PDF, DOCX, and PPT/PPTX files.

Returns
-------
parse_file() → (text: str, images: list[tuple[bytes, str]])
    text   : all extracted text joined into a single string
    images : list of (raw_image_bytes, location_hint) tuples,
             e.g. (b'\x89PNG...', 'PDF page 2, image 1')
"""
import os


# ── Public API ─────────────────────────────────────────────────────────────────

def parse_file(file_path: str, file_type: str) -> tuple:
    """
    Parse a file and return extracted text and embedded images.

    Args:
        file_path : Absolute path to the uploaded file
        file_type : File extension without dot (pdf, docx, pptx, ppt)

    Returns:
        (text, images) where text is a str and images is a list of
        (image_bytes: bytes, location_hint: str) tuples.
    """
    file_type = file_type.lower()

    if file_type == "pdf":
        return _parse_pdf(file_path)
    elif file_type == "docx":
        return _parse_docx(file_path)
    elif file_type in ("ppt", "pptx"):
        return _parse_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file type: {file_type}")


# ── PDF ────────────────────────────────────────────────────────────────────────

def _parse_pdf(file_path: str) -> tuple:
    """
    Extract text and images from a PDF.

    Primary  : PyMuPDF (fitz) — best text layout + full image extraction
    Fallback : pdfplumber → PyPDF2 (text only, no images)
    """
    text   = ""
    images = []

    # ── Attempt 1: PyMuPDF ────────────────────────────────────────────────────
    try:
        import fitz  # PyMuPDF

        doc = fitz.open(file_path)
        for page_num, page in enumerate(doc, 1):

            # Text
            page_text = page.get_text()
            if page_text.strip():
                text += page_text + "\n\n"

            # Images
            for img_idx, img_info in enumerate(page.get_images(full=True), 1):
                xref = img_info[0]
                try:
                    base_image = doc.extract_image(xref)
                    img_bytes  = base_image.get("image", b"")
                    # Skip tiny images that are likely decorators/icons (< 2 KB)
                    if img_bytes and len(img_bytes) >= 2048:
                        images.append(
                            (img_bytes, f"PDF page {page_num}, image {img_idx}")
                        )
                except Exception:
                    continue

        doc.close()
        if text.strip():
            return text.strip(), images

    except ImportError:
        pass  # PyMuPDF not installed — fall through
    except Exception:
        pass

    # ── Attempt 2: pdfplumber (text only) ────────────────────────────────────
    try:
        import pdfplumber

        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n\n"
        if text.strip():
            return text.strip(), images  # images list may be empty here

    except Exception:
        pass

    # ── Attempt 3: PyPDF2 (text only) ────────────────────────────────────────
    try:
        import PyPDF2

        with open(file_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n\n"
        if text.strip():
            return text.strip(), images

    except Exception:
        pass

    raise RuntimeError(
        "Could not extract text from PDF. "
        "The file may be scanned or entirely image-based."
    )


# ── DOCX ───────────────────────────────────────────────────────────────────────

def _parse_docx(file_path: str) -> tuple:
    """Extract text (paragraphs + tables) and embedded images from a DOCX file."""
    try:
        from docx import Document

        doc        = Document(file_path)
        paragraphs = []
        images     = []

        # Paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text.strip())

        # Tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(
                    cell.text.strip() for cell in row.cells if cell.text.strip()
                )
                if row_text:
                    paragraphs.append(row_text)

        # Images — stored as relationship targets in the document package
        img_count = 0
        for rel in doc.part.rels.values():
            if "image" not in rel.reltype.lower():
                continue
            try:
                img_bytes = rel.target_part.blob
                if img_bytes and len(img_bytes) >= 2048:
                    img_count += 1
                    images.append(
                        (img_bytes, f"DOCX embedded image {img_count}")
                    )
            except Exception:
                continue

        return "\n\n".join(paragraphs), images

    except Exception as exc:
        raise RuntimeError(f"Could not extract content from DOCX: {exc}")


# ── PPTX ───────────────────────────────────────────────────────────────────────

def _parse_pptx(file_path: str) -> tuple:
    """Extract text and images from every slide in a PPT/PPTX file."""
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        prs        = Presentation(file_path)
        slides_txt = []
        images     = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_content = [f"--- Slide {slide_num} ---"]

            for shape in slide.shapes:
                # Text
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text.strip())

                # Picture shapes
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        img_bytes = shape.image.blob
                        if img_bytes and len(img_bytes) >= 2048:
                            images.append(
                                (img_bytes, f"Slide {slide_num}")
                            )
                    except Exception:
                        continue

            # Only include slides that had some content
            if len(slide_content) > 1:
                slides_txt.append("\n".join(slide_content))

        return "\n\n".join(slides_txt), images

    except Exception as exc:
        raise RuntimeError(f"Could not extract content from PPT/PPTX: {exc}")
