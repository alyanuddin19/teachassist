"""Parse PDF, DOCX, PPT/PPTX, TXT, and LaTeX files into text plus embedded images."""


def parse_file(file_path: str, file_type: str) -> tuple[str, list[tuple[bytes, str]]]:
    file_type = file_type.lower()
    if file_type == "pdf":
        return _parse_pdf(file_path)
    if file_type == "docx":
        return _parse_docx(file_path)
    if file_type in ("ppt", "pptx"):
        return _parse_pptx(file_path)
    if file_type in ("txt", "tex", "latex"):
        return _parse_text_file(file_path)
    raise ValueError(f"Unsupported file type: {file_type}")


def _parse_pdf(file_path: str) -> tuple[str, list[tuple[bytes, str]]]:
    text = ""
    images: list[tuple[bytes, str]] = []

    try:
        import fitz

        doc = fitz.open(file_path)
        for page_number, page in enumerate(doc, 1):
            page_text = page.get_text()
            if page_text.strip():
                text += page_text + "\n\n"

            for image_index, image_info in enumerate(page.get_images(full=True), 1):
                xref = image_info[0]
                try:
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image.get("image", b"")
                    if image_bytes and len(image_bytes) >= 2048:
                        images.append((image_bytes, f"PDF page {page_number}, image {image_index}"))
                except Exception:
                    continue
        doc.close()
        if text.strip():
            return text.strip(), images
    except Exception:
        pass

    try:
        import pdfplumber

        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n\n"
        if text.strip():
            return text.strip(), images
    except Exception:
        pass

    try:
        import PyPDF2

        with open(file_path, "rb") as file_handle:
            reader = PyPDF2.PdfReader(file_handle)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n\n"
        if text.strip():
            return text.strip(), images
    except Exception:
        pass

    raise RuntimeError("Could not extract text from PDF. The file may be scanned or image-based.")


def _parse_docx(file_path: str) -> tuple[str, list[tuple[bytes, str]]]:
    from docx import Document

    doc = Document(file_path)
    paragraphs: list[str] = []
    images: list[tuple[bytes, str]] = []

    for para in doc.paragraphs:
        if para.text.strip():
            paragraphs.append(para.text.strip())

    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                paragraphs.append(row_text)

    image_count = 0
    for rel in doc.part.rels.values():
        if "image" not in rel.reltype.lower():
            continue
        try:
            image_bytes = rel.target_part.blob
            if image_bytes and len(image_bytes) >= 2048:
                image_count += 1
                images.append((image_bytes, f"DOCX embedded image {image_count}"))
        except Exception:
            continue

    return "\n\n".join(paragraphs), images


def _parse_pptx(file_path: str) -> tuple[str, list[tuple[bytes, str]]]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    presentation = Presentation(file_path)
    slide_text: list[str] = []
    images: list[tuple[bytes, str]] = []

    for slide_number, slide in enumerate(presentation.slides, 1):
        content = [f"--- Slide {slide_number} ---"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                content.append(shape.text.strip())

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image_bytes = shape.image.blob
                    if image_bytes and len(image_bytes) >= 2048:
                        images.append((image_bytes, f"Slide {slide_number}"))
                except Exception:
                    continue

        if len(content) > 1:
            slide_text.append("\n".join(content))

    return "\n\n".join(slide_text), images


def _parse_text_file(file_path: str) -> tuple[str, list[tuple[bytes, str]]]:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as file_handle:
        return file_handle.read(), []
