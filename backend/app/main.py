from fastapi import FastAPI, Depends, HTTPException, UploadFile, File, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse, FileResponse
from sqlalchemy.orm import Session
from sqlalchemy import or_, func, text
from dotenv import load_dotenv
from io import BytesIO
from datetime import datetime, timezone, timedelta
from pathlib import Path
from uuid import uuid4
import re
import os
import json
import random
import smtplib
import pandas as pd
from email.message import EmailMessage

from groq import Groq
from pypdf import PdfReader
from docx import Document
from pptx import Presentation

from .database import Base, engine, SessionLocal
from . import models, schemas
from .auth import hash_password
from app.core.rag_engine import RAGEngine
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.units import inch
from typing import List, Dict, Any
from collections import defaultdict
from reportlab.platypus import SimpleDocTemplate, Spacer, Preformatted


from reportlab.platypus import SimpleDocTemplate, Paragraph
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.pagesizes import A4

from app.routers import gap_analysis_clean

TRANSFORM_MARKSHEET_RETENTION_DAYS = 28
from app.routers import prompt_generator
from app.routers import transformation
from app.t2_transform_loader import load_t2_transform_app, load_t2_excel_builder



load_dotenv()
UPLOAD_ROOT = Path(__file__).resolve().parents[1] / "uploads"
TASK_SUBMISSION_DIR = UPLOAD_ROOT / "task_submissions"
TASK_SUBMISSION_DIR.mkdir(parents=True, exist_ok=True)


def parse_allowed_origins() -> list[str]:
    raw = os.getenv("ALLOWED_ORIGINS", "").strip()
    if raw:
        return [origin.strip() for origin in raw.split(",") if origin.strip()]

    return [
        "http://localhost:4200",
        "http://127.0.0.1:4200",
    ]

# ------------------ APP & DB ------------------
Base.metadata.create_all(bind=engine)
app = FastAPI(title="TeachAssist Backend")

app.include_router(gap_analysis_clean.router)
app.include_router(prompt_generator.router)
app.include_router(transformation.router)

t2_transform_app = load_t2_transform_app()
if t2_transform_app is not None:
    app.mount("/transform-tool", t2_transform_app)

# ------------------ CORS ------------------
app.add_middleware(
    CORSMiddleware,
    allow_origins=parse_allowed_origins(),
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


@app.get("/health")
def healthcheck():
    return {"status": "ok"}
# ------------------ DB DEP ------------------
def get_db():
    db = SessionLocal()
    try:
        yield db
    finally:
        db.close()


def normalize_text(value: str | None, uppercase: bool = False) -> str:
    cleaned = (value or "").strip()
    return cleaned.upper() if uppercase else cleaned


def slugify_name(name: str) -> str:
    slug = re.sub(r"[^a-z0-9]+", ".", name.lower()).strip(".")
    return slug or "user"


def generate_unique_username(name: str, model, field_name: str, db: Session) -> str:
    base = slugify_name(name)
    candidate = base
    suffix = 1

    while db.query(model).filter(getattr(model, field_name) == candidate).first():
        candidate = f"{base}{suffix}"
        suffix += 1

    return candidate


def generate_easy_password(name: str, prefix: str = "edu") -> str:
    clean = re.sub(r"[^A-Za-z]", "", name).lower()
    seed = (clean[:5] or prefix)
    return f"{seed}{random.randint(100, 999)}"


def generate_unique_uid(prefix: str, model, field_name: str, db: Session) -> str:
    while True:
        uid = f"{prefix}-{uuid4().hex[:8].upper()}"
        if not db.query(model).filter(getattr(model, field_name) == uid).first():
            return uid


def build_student_email(roll_no: str) -> str:
    return f"{normalize_text(roll_no, uppercase=True)}@students.uitu.edu.pk"


def build_student_batch_table_name(department: str, batch: str) -> str:
    department_slug = re.sub(r"[^a-z0-9]+", "_", (department or "general").strip().lower()).strip("_") or "general"
    batch_slug = re.sub(r"[^a-z0-9]+", "_", (batch or "batch").strip().lower()).strip("_") or "batch"
    return f"students_{department_slug}_{batch_slug}"


def ensure_student_batch_table(table_name: str, db: Session) -> None:
    db.execute(text(f"""
        CREATE TABLE IF NOT EXISTS {table_name} (
            LIKE students INCLUDING DEFAULTS INCLUDING CONSTRAINTS
        )
    """))
    db.execute(text(f"CREATE UNIQUE INDEX IF NOT EXISTS {table_name}_roll_no_idx ON {table_name} (roll_no)"))
    db.execute(text(f"CREATE UNIQUE INDEX IF NOT EXISTS {table_name}_username_idx ON {table_name} (username)"))
    db.commit()


def mirror_student_to_batch_table(student: models.Student, table_name: str, db: Session) -> None:
    db.execute(
        text(f"""
            INSERT INTO {table_name} (
                id, student_name, roll_no, username, password, contact_no, email,
                program, teacher_id, semester, section, batch, department, created_at
            )
            VALUES (
                :id, :student_name, :roll_no, :username, :password, :contact_no, :email,
                :program, :teacher_id, :semester, :section, :batch, :department, :created_at
            )
            ON CONFLICT (roll_no)
            DO UPDATE SET
                student_name = EXCLUDED.student_name,
                username = EXCLUDED.username,
                password = EXCLUDED.password,
                contact_no = EXCLUDED.contact_no,
                email = EXCLUDED.email,
                program = EXCLUDED.program,
                teacher_id = EXCLUDED.teacher_id,
                semester = EXCLUDED.semester,
                section = EXCLUDED.section,
                batch = EXCLUDED.batch,
                department = EXCLUDED.department
        """),
        {
            "id": student.id,
            "student_name": student.student_name,
            "roll_no": student.roll_no,
            "username": student.username,
            "password": student.password,
            "contact_no": student.contact_no,
            "email": student.email,
            "program": student.program,
            "teacher_id": student.teacher_id,
            "semester": student.semester,
            "section": student.section,
            "batch": student.batch,
            "department": student.department,
            "created_at": student.created_at,
        }
    )


def parse_student_import_rows(file_bytes: bytes) -> list[dict]:
    dataframe = pd.read_excel(BytesIO(file_bytes), engine="openpyxl")
    normalized_columns = {
        str(column).strip().lower().replace(" ", "_").replace("-", "_"): column
        for column in dataframe.columns
    }

    def read_value(row, *candidates: str):
        for candidate in candidates:
            original = normalized_columns.get(candidate)
            if original is not None:
                value = row.get(original)
                if pd.notna(value):
                    return str(value).strip()
        return ""

    parsed_rows: list[dict] = []
    for _, row in dataframe.iterrows():
        student_name = read_value(row, "student_name", "full_name", "name")
        roll_no = read_value(row, "roll_no", "roll_number", "student_id", "id")
        if not student_name or not roll_no:
            continue

        parsed_rows.append({
            "student_name": student_name,
            "roll_no": roll_no.upper(),
            "contact_no": read_value(row, "contact_no", "contact", "contact_number", "phone"),
            "email": read_value(row, "email", "email_address"),
            "program": read_value(row, "program", "programs", "programme"),
            "department": read_value(row, "department", "dept").upper(),
            "batch": read_value(row, "batch").upper(),
            "section": read_value(row, "section").upper(),
            "semester": read_value(row, "semester"),
        })
    return parsed_rows


def send_credentials_email(recipient: str, subject: str, body: str) -> bool:
    host = os.getenv("SMTP_HOST")
    port = os.getenv("SMTP_PORT")
    username = os.getenv("SMTP_USERNAME")
    password = os.getenv("SMTP_PASSWORD")
    sender = os.getenv("SMTP_FROM_EMAIL") or username

    if not all([host, port, sender]):
        return False

    message = EmailMessage()
    message["From"] = sender
    message["To"] = recipient
    message["Subject"] = subject
    message.set_content(body)

    try:
        with smtplib.SMTP(host, int(port)) as server:
            server.starttls()
            if username and password:
                server.login(username, password)
            server.send_message(message)
        return True
    except Exception:
        return False


def build_teacher_record_payload(record: models.TeacherCourse, course: models.Course | None, students: list[models.Student]) -> dict:
    return {
        "id": record.id,
        "semester": record.semester,
        "section": record.section or "",
        "batch": record.batch or "",
        "department": record.department or "",
        "threshold_percentage": record.threshold_percentage if record.threshold_percentage is not None else 50,
        "course_code": course.course_code if course else "",
        "course_name": course.course_name if course else "",
        "students": [
            {
                "id": student.id,
                "student_name": student.student_name,
                "roll_no": student.roll_no or "",
                "semester": student.semester,
                "section": student.section or "",
                "batch": student.batch or "",
                "department": student.department or ""
            }
            for student in students
        ]
    }


def build_student_task_payload(task: models.StudentTask, teacher: models.Teacher | None = None) -> dict:
    attachment_url = (
        f"/api/student/tasks/{task.id}/attachment"
        if task.answer_attachment_path
        else None
    )
    task_attachment_url = (
        f"/api/student/tasks/{task.id}/reference"
        if task.task_attachment_path
        else None
    )
    return {
        "id": task.id,
        "title": task.title,
        "description": task.description or "",
        "question_content": task.question_content or "",
        "status": task.status,
        "source_module": task.source_module,
        "teacher_name": teacher.teacher_name if teacher else "",
        "assigned_roll_no": task.assigned_roll_no or "",
        "clo": task.clo or "",
        "course_name": task.course_name_snapshot or "",
        "task_attachment_name": task.task_attachment_name or "",
        "task_attachment_url": task_attachment_url,
        "answer_text": task.answer_text or "",
        "answer_attachment_name": task.answer_attachment_name or "",
        "answer_attachment_url": attachment_url,
        "submitted_at": task.submitted_at.isoformat() if task.submitted_at else None,
        "reviewed_at": task.reviewed_at.isoformat() if task.reviewed_at else None,
        "teacher_decision": task.teacher_decision or "",
        "teacher_feedback": task.teacher_feedback or "",
        "teacher_score": task.teacher_score or "",
        "created_at": task.created_at.isoformat() if task.created_at else None
    }


def build_teacher_submission_payload(
    task: models.StudentTask,
    student: models.Student | None,
    course: models.Course | None,
    teacher: models.Teacher | None = None
) -> dict:
    return {
        "task_id": task.id,
        "title": task.title,
        "status": task.status,
        "student_id": student.id if student else None,
        "student_name": student.student_name if student else "",
        "student_roll_no": student.roll_no if student else "",
        "teacher_name": teacher.teacher_name if teacher else "",
        "course_code": course.course_code if course else "",
        "course_name": task.course_name_snapshot or (course.course_name if course else ""),
        "question_content": task.question_content or "",
        "task_attachment_name": task.task_attachment_name or "",
        "task_attachment_url": f"/api/student/tasks/{task.id}/reference" if task.task_attachment_path else "",
        "answer_text": task.answer_text or "",
        "answer_attachment_name": task.answer_attachment_name or "",
        "answer_attachment_url": f"/api/student/tasks/{task.id}/attachment" if task.answer_attachment_path else "",
        "teacher_feedback": task.teacher_feedback or "",
        "teacher_score": task.teacher_score or "",
        "teacher_decision": task.teacher_decision or "",
        "submitted_at": task.submitted_at.isoformat() if task.submitted_at else None,
        "reviewed_at": task.reviewed_at.isoformat() if task.reviewed_at else None,
        "created_at": task.created_at.isoformat() if task.created_at else None,
        "clo": task.clo or "",
        "source_module": task.source_module
    }


def get_students_for_teacher_record(record: models.TeacherCourse, db: Session) -> list[models.Student]:
    assignment_rows = db.query(models.TeacherStudentAssignment).filter(
        models.TeacherStudentAssignment.teacher_id == record.teacher_id,
        models.TeacherStudentAssignment.teacher_course_id == record.id
    ).all()

    if not assignment_rows:
        return []

    student_ids = [row.student_id for row in assignment_rows]
    return db.query(models.Student).filter(
        models.Student.id.in_(student_ids)
    ).order_by(models.Student.student_name.asc()).all()


def get_teacher_profile_data(teacher: models.Teacher, db: Session) -> dict:
    records = db.query(models.TeacherCourse).filter(
        models.TeacherCourse.teacher_id == teacher.id
    ).order_by(models.TeacherCourse.created_at.desc(), models.TeacherCourse.id.desc()).all()

    course_map: dict[int, models.Course] = {
        course.id: course
        for course in db.query(models.Course).all()
    }

    payload_records: list[dict] = []
    for record in records:
        students = get_students_for_teacher_record(record, db)
        payload_records.append(
            build_teacher_record_payload(record, course_map.get(record.course_id), students)
        )

    return {
        "teacher_id": teacher.id,
        "teacher_uid": teacher.teacher_uid or "",
        "teacher_name": teacher.teacher_name,
        "department": teacher.department or "",
        "username": teacher.username or "",
        "setup_complete": len(payload_records) > 0,
        "records": payload_records
    }


def assign_matching_students_to_teacher(
    teacher_id: int,
    teacher_course_id: int,
    course_id: int,
    semester: int,
    batch: str,
    department: str,
    section: str | None,
    db: Session
) -> int:
    student_query = db.query(models.Student).filter(
        models.Student.semester == semester,
        models.Student.batch == batch,
        models.Student.department == department
    )

    if section:
        student_query = student_query.filter(models.Student.section == section)

    matched_students = student_query.all()
    for student in matched_students:
        exists = db.query(models.TeacherStudentAssignment).filter(
            models.TeacherStudentAssignment.teacher_id == teacher_id,
            models.TeacherStudentAssignment.student_id == student.id,
            models.TeacherStudentAssignment.teacher_course_id == teacher_course_id
        ).first()
        if exists:
            continue

        db.add(models.TeacherStudentAssignment(
            teacher_id=teacher_id,
            student_id=student.id,
            teacher_course_id=teacher_course_id,
            course_id=course_id,
            semester=semester,
            section=section or None,
            batch=batch,
            department=department
        ))

    db.commit()
    return len(matched_students)


def sync_teacher_record_students(
    record: models.TeacherCourse,
    course_id: int,
    semester: int,
    batch: str,
    department: str,
    section: str | None,
    db: Session
) -> int:
    db.query(models.TeacherStudentAssignment).filter(
        models.TeacherStudentAssignment.teacher_id == record.teacher_id,
        models.TeacherStudentAssignment.teacher_course_id == record.id
    ).delete()
    db.commit()

    return assign_matching_students_to_teacher(
        record.teacher_id,
        record.id,
        course_id,
        semester,
        batch,
        department,
        section,
        db
    )


def create_teacher_record(data: dict, db: Session) -> dict:
    teacher_name = normalize_text(data.get("teacher_name"))
    semester = data.get("semester")
    section = normalize_text(data.get("section"), uppercase=True)
    course_code = normalize_text(data.get("course_code"), uppercase=True)
    batch = normalize_text(data.get("batch"), uppercase=True)
    department = normalize_text(data.get("department"), uppercase=True)
    threshold_input = data.get("threshold_percentage")

    if not all([teacher_name, semester, batch, department, course_code]):
        raise HTTPException(status_code=400, detail="Missing required fields")

    try:
        threshold_percentage = int(float(threshold_input if threshold_input is not None else 50))
    except (TypeError, ValueError):
        raise HTTPException(status_code=400, detail="Threshold must be a valid percentage")

    if threshold_percentage < 1 or threshold_percentage > 100:
        raise HTTPException(status_code=400, detail="Threshold must be between 1 and 100")

    teacher = db.query(models.Teacher).filter(
        models.Teacher.teacher_name == teacher_name
    ).first()
    if not teacher:
        raise HTTPException(status_code=404, detail="Teacher not found")

    course = db.query(models.Course).filter(
        models.Course.course_code == course_code
    ).first()
    if not course:
        raise HTTPException(status_code=404, detail="Course not found")

    semester_value = int(semester)

    existing = db.query(models.TeacherCourse).filter(
        models.TeacherCourse.teacher_id == teacher.id,
        models.TeacherCourse.course_id == course.id,
        models.TeacherCourse.semester == semester_value,
        models.TeacherCourse.batch == batch,
        models.TeacherCourse.department == department,
        models.TeacherCourse.section == (section or None)
    ).first()

    if existing:
        existing.threshold_percentage = threshold_percentage
        db.commit()
        db.refresh(existing)
        assigned_count = sync_teacher_record_students(
            existing,
            course.id,
            semester_value,
            batch,
            department,
            section or None,
            db
        )
        return {
            "status": "already_exists",
            "assigned_students": assigned_count,
            "record": build_teacher_record_payload(
                existing,
                course,
                get_students_for_teacher_record(existing, db)
            )
        }

    record = models.TeacherCourse(
        teacher_id=teacher.id,
        course_id=course.id,
        semester=semester_value,
        section=section or None,
        batch=batch,
        department=department,
        threshold_percentage=threshold_percentage
    )
    db.add(record)
    db.commit()
    db.refresh(record)

    assigned_count = sync_teacher_record_students(
        record,
        course.id,
        semester_value,
        batch,
        department,
        section or None,
        db
    )

    return {
        "status": "record_added",
        "assigned_students": assigned_count,
        "record": build_teacher_record_payload(
            record,
            course,
            get_students_for_teacher_record(record, db)
        )
    }


def update_teacher_record_threshold(record_id: int, data: dict, db: Session) -> dict:
    teacher_name = normalize_text(data.get("teacher_name"))
    threshold_input = data.get("threshold_percentage")

    if not teacher_name:
        raise HTTPException(status_code=400, detail="Teacher name is required")

    try:
        threshold_percentage = int(float(threshold_input))
    except (TypeError, ValueError):
        raise HTTPException(status_code=400, detail="Threshold must be a valid percentage")

    if threshold_percentage < 1 or threshold_percentage > 100:
        raise HTTPException(status_code=400, detail="Threshold must be between 1 and 100")

    teacher = db.query(models.Teacher).filter(
        models.Teacher.teacher_name == teacher_name
    ).first()
    if not teacher:
        raise HTTPException(status_code=404, detail="Teacher not found")

    record = db.query(models.TeacherCourse).filter(
        models.TeacherCourse.id == record_id,
        models.TeacherCourse.teacher_id == teacher.id
    ).first()
    if not record:
        raise HTTPException(status_code=404, detail="Teaching record not found")

    record.threshold_percentage = threshold_percentage
    db.commit()
    db.refresh(record)

    course = db.query(models.Course).filter(models.Course.id == record.course_id).first()
    return {
        "status": "updated",
        "record": build_teacher_record_payload(
            record,
            course,
            get_students_for_teacher_record(record, db)
        )
    }


def update_teacher_record(record_id: int, data: dict, db: Session) -> dict:
    teacher_name = normalize_text(data.get("teacher_name"))
    semester = data.get("semester")
    section = normalize_text(data.get("section"), uppercase=True)
    course_code = normalize_text(data.get("course_code"), uppercase=True)
    batch = normalize_text(data.get("batch"), uppercase=True)
    department = normalize_text(data.get("department"), uppercase=True)
    threshold_input = data.get("threshold_percentage")

    if not all([teacher_name, semester, batch, department, course_code]):
        raise HTTPException(status_code=400, detail="Missing required fields")

    try:
        semester_value = int(semester)
    except (TypeError, ValueError):
        raise HTTPException(status_code=400, detail="Semester must be a valid number")

    try:
        threshold_percentage = int(float(threshold_input if threshold_input is not None else 50))
    except (TypeError, ValueError):
        raise HTTPException(status_code=400, detail="Threshold must be a valid percentage")

    if threshold_percentage < 1 or threshold_percentage > 100:
        raise HTTPException(status_code=400, detail="Threshold must be between 1 and 100")

    teacher = db.query(models.Teacher).filter(
        models.Teacher.teacher_name == teacher_name
    ).first()
    if not teacher:
        raise HTTPException(status_code=404, detail="Teacher not found")

    record = db.query(models.TeacherCourse).filter(
        models.TeacherCourse.id == record_id,
        models.TeacherCourse.teacher_id == teacher.id
    ).first()
    if not record:
        raise HTTPException(status_code=404, detail="Teaching record not found")

    course = db.query(models.Course).filter(
        models.Course.course_code == course_code
    ).first()
    if not course:
        raise HTTPException(status_code=404, detail="Course not found")

    duplicate = db.query(models.TeacherCourse).filter(
        models.TeacherCourse.teacher_id == teacher.id,
        models.TeacherCourse.course_id == course.id,
        models.TeacherCourse.semester == semester_value,
        models.TeacherCourse.batch == batch,
        models.TeacherCourse.department == department,
        models.TeacherCourse.section == (section or None),
        models.TeacherCourse.id != record.id
    ).first()
    if duplicate:
        raise HTTPException(status_code=400, detail="Another teaching record already exists with the same course and class details")

    record.course_id = course.id
    record.semester = semester_value
    record.section = section or None
    record.batch = batch
    record.department = department
    record.threshold_percentage = threshold_percentage
    db.commit()
    db.refresh(record)

    assigned_count = sync_teacher_record_students(
        record,
        course.id,
        semester_value,
        batch,
        department,
        section or None,
        db
    )

    return {
        "status": "updated",
        "assigned_students": assigned_count,
        "record": build_teacher_record_payload(
            record,
            course,
            get_students_for_teacher_record(record, db)
        )
    }


def delete_teacher_record(record_id: int, teacher_name: str, db: Session) -> dict:
    normalized_teacher_name = normalize_text(teacher_name)
    if not normalized_teacher_name:
        raise HTTPException(status_code=400, detail="Teacher name is required")

    teacher = db.query(models.Teacher).filter(
        models.Teacher.teacher_name == normalized_teacher_name
    ).first()
    if not teacher:
        raise HTTPException(status_code=404, detail="Teacher not found")

    record = db.query(models.TeacherCourse).filter(
        models.TeacherCourse.id == record_id,
        models.TeacherCourse.teacher_id == teacher.id
    ).first()
    if not record:
        raise HTTPException(status_code=404, detail="Teaching record not found")

    db.query(models.TeacherStudentAssignment).filter(
        models.TeacherStudentAssignment.teacher_id == teacher.id,
        models.TeacherStudentAssignment.teacher_course_id == record.id
    ).delete()

    db.delete(record)
    db.commit()

    return {"status": "deleted", "record_id": record_id}
# =========================================================
# LOGIN API
# =========================================================
@app.post("/api/login")
def login(data: dict, db: Session = Depends(get_db)):
    teacher_login = data.get("teacher_name")
    password = data.get("password")

    if not teacher_login or not password:
        raise HTTPException(status_code=400, detail="Missing credentials")

    teacher = db.query(models.Teacher).filter(
        or_(
            models.Teacher.username == teacher_login,
            models.Teacher.teacher_name == teacher_login
        )
    ).first()

    if not teacher:
        return {"status": "not_found"}

    if teacher.password == "PENDING_SETUP":
        return {"status": "pending_setup"}

    if teacher.password != password:
        return {"status": "invalid_password"}

    return {
        "status": "found",
        "teacher_id": teacher.id,
        "teacher_name": teacher.teacher_name,
        "teacher_uid": teacher.teacher_uid or "",
        "department": teacher.department or "",
        "username": teacher.username or ""
    }


@app.post("/api/admin/login")
def admin_login(data: dict):
    username = normalize_text(data.get("username"))
    password = normalize_text(data.get("password"))

    if not username or not password:
        raise HTTPException(status_code=400, detail="Missing credentials")

    if username not in {"amdin", "admin"} or password != "admin123":
        return {"status": "invalid_credentials"}

    return {
        "status": "found",
        "admin_name": "System Admin"
    }


@app.post("/api/hod/login")
def hod_login(data: dict, db: Session = Depends(get_db)):
    username = normalize_text(data.get("username"))
    password = data.get("password")

    if not username or not password:
        raise HTTPException(status_code=400, detail="Missing credentials")

    hod = db.query(models.HeadOfDepartment).filter(
        models.HeadOfDepartment.username == username
    ).first()

    if not hod:
        return {"status": "not_found"}

    if hod.password != password:
        return {"status": "invalid_password"}

    return {
        "status": "found",
        "hod_id": hod.id,
        "hod_name": hod.full_name,
        "hod_uid": hod.hod_uid,
        "department": hod.department,
        "username": hod.username
    }


@app.post("/api/student/login")
def student_login(data: dict, db: Session = Depends(get_db)):
    student_code = data.get("student_code")
    password = data.get("password")

    if not student_code or not password:
        raise HTTPException(status_code=400, detail="Missing credentials")

    student = db.query(models.Student).filter(
        models.Student.roll_no == student_code
    ).first()

    if not student:
        return {"status": "not_found"}

    if student.password != password:
        return {"status": "invalid_password"}

    return {
        "status": "found",
        "student_id": student.id,
        "student_code": student.roll_no or student.student_name,
        "student_name": student.student_name,
        "username": student.username or ""
    }


@app.post("/api/student/change-password")
def change_student_password(data: dict, db: Session = Depends(get_db)):
    student_id = data.get("student_id")
    current_password = data.get("current_password") or ""
    new_password = data.get("new_password") or ""

    if not student_id or not current_password or not new_password:
        raise HTTPException(status_code=400, detail="Student ID, current password, and new password are required")

    if len(new_password) < 6:
        raise HTTPException(status_code=400, detail="New password must be at least 6 characters")

    student = db.query(models.Student).filter(models.Student.id == student_id).first()
    if not student:
        raise HTTPException(status_code=404, detail="Student not found")

    if student.password != current_password:
        raise HTTPException(status_code=400, detail="Current password is incorrect")

    student.password = new_password
    db.commit()

    return {"status": "password_changed"}


@app.get("/api/student/tasks")
def get_student_tasks(student_id: int, db: Session = Depends(get_db)):
    student = db.query(models.Student).filter(
        models.Student.id == student_id
    ).first()

    if not student:
        raise HTTPException(status_code=404, detail="Student not found")

    tasks = db.query(models.StudentTask).filter(
        models.StudentTask.student_id == student.id
    ).order_by(models.StudentTask.created_at.desc()).all()
    teacher_ids = [task.teacher_id for task in tasks if task.teacher_id]
    teacher_map = {
        teacher.id: teacher
        for teacher in db.query(models.Teacher).filter(models.Teacher.id.in_(teacher_ids)).all()
    } if teacher_ids else {}

    return {
        "student_id": student.id,
        "student_code": student.roll_no or student.student_name,
        "student_name": student.student_name,
        "tasks": [build_student_task_payload(task, teacher_map.get(task.teacher_id)) for task in tasks]
    }


@app.delete("/api/student/tasks/{task_id}")
def delete_student_task(task_id: int, student_id: int, db: Session = Depends(get_db)):
    task = db.query(models.StudentTask).filter(
        models.StudentTask.id == task_id,
        models.StudentTask.student_id == student_id
    ).first()

    if not task:
        raise HTTPException(status_code=404, detail="Task not found")

    for file_path in [task.task_attachment_path, task.answer_attachment_path]:
        if file_path:
            try:
                path = Path(file_path)
                if path.exists():
                    path.unlink()
            except Exception:
                pass

    db.delete(task)
    db.commit()

    return {"status": "deleted", "task_id": task_id}


@app.get("/api/student/tasks/{task_id}/attachment")
def get_task_attachment(task_id: int, db: Session = Depends(get_db)):
    task = db.query(models.StudentTask).filter(
        models.StudentTask.id == task_id
    ).first()

    if not task or not task.answer_attachment_path:
        raise HTTPException(status_code=404, detail="Attachment not found")

    file_path = Path(task.answer_attachment_path)
    if not file_path.exists():
        raise HTTPException(status_code=404, detail="Attachment file missing")

    return FileResponse(path=file_path, filename=task.answer_attachment_name or file_path.name)


@app.get("/api/student/tasks/{task_id}/reference")
def get_task_reference(task_id: int, db: Session = Depends(get_db)):
    task = db.query(models.StudentTask).filter(
        models.StudentTask.id == task_id
    ).first()

    if not task or not task.task_attachment_path:
        raise HTTPException(status_code=404, detail="Reference file not found")

    file_path = Path(task.task_attachment_path)
    if not file_path.exists():
        raise HTTPException(status_code=404, detail="Reference file missing")

    return FileResponse(path=file_path, filename=task.task_attachment_name or file_path.name)


@app.post("/api/student/tasks/assign")
async def assign_student_task(
    teacher_name: str = Form(...),
    student_roll_no: str = Form(...),
    title: str = Form(...),
    description: str = Form(""),
    question_content: str = Form(""),
    clo: str = Form(""),
    source_module: str = Form("gap_analysis"),
    course_code: str = Form(""),
    course_name: str = Form(""),
    reference_file: UploadFile | None = File(None),
    db: Session = Depends(get_db)
):
    teacher_name = normalize_text(teacher_name)
    student_roll_no = normalize_text(student_roll_no, uppercase=True)
    title = normalize_text(title)
    description = normalize_text(description)
    clo = normalize_text(clo, uppercase=True)
    source_module = normalize_text(source_module) or "gap_analysis"
    course_code = normalize_text(course_code, uppercase=True)
    course_name = normalize_text(course_name)

    if not all([teacher_name, student_roll_no, title]):
        raise HTTPException(status_code=400, detail="Missing task assignment fields")

    teacher = db.query(models.Teacher).filter(
        models.Teacher.teacher_name == teacher_name
    ).first()
    if not teacher:
        raise HTTPException(status_code=404, detail="Teacher not found")

    student = db.query(models.Student).filter(
        models.Student.roll_no == student_roll_no
    ).first()
    if not student:
        raise HTTPException(status_code=404, detail="Student with this roll number was not found")

    course = None
    if course_code:
        course = db.query(models.Course).filter(
            models.Course.course_code == course_code
        ).first()
    if not course and course_name:
        course = db.query(models.Course).filter(
            func.lower(models.Course.course_name) == course_name.lower()
        ).first()

    task = models.StudentTask(
        student_id=student.id,
        teacher_id=teacher.id,
        course_id=course.id if course else None,
        course_name_snapshot=course_name or (course.course_name if course else None),
        title=title,
        description=description or None,
        question_content=question_content or None,
        clo=clo or None,
        assigned_roll_no=student.roll_no or student_roll_no,
        source_module=source_module,
        status="assigned"
    )

    if reference_file and reference_file.filename:
        extension = Path(reference_file.filename).suffix
        saved_name = f"{uuid4().hex}{extension}"
        target_path = TASK_SUBMISSION_DIR / saved_name
        file_bytes = await reference_file.read()
        target_path.write_bytes(file_bytes)
        task.task_attachment_name = reference_file.filename
        task.task_attachment_path = str(target_path)

    db.add(task)
    db.commit()
    db.refresh(task)

    return {
        "status": "assigned",
        "task": build_student_task_payload(task, teacher),
        "student_name": student.student_name,
        "student_roll_no": student.roll_no or student_roll_no
    }


@app.post("/api/student/tasks/{task_id}/submit")
async def submit_student_task(
    task_id: int,
    student_id: int = Form(...),
    answer_text: str = Form(""),
    answer_file: UploadFile | None = File(None),
    db: Session = Depends(get_db)
):
    task = db.query(models.StudentTask).filter(
        models.StudentTask.id == task_id,
        models.StudentTask.student_id == student_id
    ).first()

    if not task:
        raise HTTPException(status_code=404, detail="Task not found")

    if task.status not in {"assigned", "returned"}:
        raise HTTPException(status_code=400, detail="This task cannot be submitted right now")

    answer_text = (answer_text or "").strip()
    if not answer_text and (not answer_file or not answer_file.filename):
        raise HTTPException(status_code=400, detail="Add a typed answer or an attachment")

    if answer_file and answer_file.filename:
        extension = Path(answer_file.filename).suffix
        saved_name = f"{uuid4().hex}{extension}"
        target_path = TASK_SUBMISSION_DIR / saved_name
        file_bytes = await answer_file.read()
        target_path.write_bytes(file_bytes)
        task.answer_attachment_name = answer_file.filename
        task.answer_attachment_path = str(target_path)

    task.answer_text = answer_text or None
    task.submitted_at = datetime.now(timezone.utc)
    task.status = "submitted"
    task.teacher_decision = None
    task.teacher_feedback = None
    task.teacher_score = None
    task.reviewed_at = None

    db.commit()
    db.refresh(task)

    teacher = db.query(models.Teacher).filter(models.Teacher.id == task.teacher_id).first() if task.teacher_id else None
    return {"status": "submitted", "task": build_student_task_payload(task, teacher)}


@app.get("/api/teacher/notifications")
def get_teacher_notifications(teacher_name: str, db: Session = Depends(get_db)):
    teacher = db.query(models.Teacher).filter(
        models.Teacher.teacher_name == teacher_name
    ).first()

    if not teacher:
        raise HTTPException(status_code=404, detail="Teacher not found")

    tasks = db.query(models.StudentTask).filter(
        models.StudentTask.teacher_id == teacher.id,
        models.StudentTask.status.in_(["submitted", "passed", "failed"])
    ).order_by(
        models.StudentTask.submitted_at.desc(),
        models.StudentTask.created_at.desc()
    ).all()

    course_ids = [task.course_id for task in tasks if task.course_id]
    student_ids = [task.student_id for task in tasks]
    teacher_ids = [task.teacher_id for task in tasks if task.teacher_id]
    course_map = {
        course.id: course for course in db.query(models.Course).filter(models.Course.id.in_(course_ids)).all()
    } if course_ids else {}
    student_map = {
        student.id: student for student in db.query(models.Student).filter(models.Student.id.in_(student_ids)).all()
    } if student_ids else {}
    teacher_map = {
        teacher.id: teacher for teacher in db.query(models.Teacher).filter(models.Teacher.id.in_(teacher_ids)).all()
    } if teacher_ids else {}

    notifications = [
        build_teacher_submission_payload(task, student_map.get(task.student_id), course_map.get(task.course_id), teacher_map.get(task.teacher_id))
        for task in tasks
    ]
    pending_count = len([item for item in notifications if item["status"] == "submitted"])

    return {
        "teacher_id": teacher.id,
        "pending_count": pending_count,
        "submissions": notifications
    }


@app.post("/api/teacher/tasks/{task_id}/review")
def review_student_task(task_id: int, data: dict, db: Session = Depends(get_db)):
    teacher_name = normalize_text(data.get("teacher_name"))
    decision = normalize_text(data.get("decision"), uppercase=True)
    feedback = (data.get("feedback") or "").strip()
    score = (data.get("score") or "").strip()

    if decision not in {"PASS", "FAIL"}:
        raise HTTPException(status_code=400, detail="Decision must be PASS or FAIL")

    teacher = db.query(models.Teacher).filter(
        models.Teacher.teacher_name == teacher_name
    ).first()
    if not teacher:
        raise HTTPException(status_code=404, detail="Teacher not found")

    task = db.query(models.StudentTask).filter(
        models.StudentTask.id == task_id,
        models.StudentTask.teacher_id == teacher.id
    ).first()
    if not task:
        raise HTTPException(status_code=404, detail="Task not found")

    cleanup_expired_transform_marksheets(db, teacher.id)

    task.teacher_decision = decision
    task.teacher_feedback = feedback or None
    task.teacher_score = score or None
    task.reviewed_at = datetime.now(timezone.utc)
    task.status = "passed" if decision == "PASS" else "failed"

    db.commit()
    db.refresh(task)

    auto_updated_sheet_id = None
    auto_updated_sheet_download_url = ""
    auto_applied_updates = 0
    latest_marksheet = None
    if task.course_id:
        latest_marksheet = db.query(models.TransformMarksheet).filter(
            models.TransformMarksheet.teacher_id == teacher.id,
            models.TransformMarksheet.course_id == task.course_id
        ).order_by(models.TransformMarksheet.created_at.desc(), models.TransformMarksheet.id.desc()).first()
    elif task.course_name_snapshot:
        normalized_course_name = normalize_text(task.course_name_snapshot)
        teacher_marksheets = db.query(models.TransformMarksheet).filter(
            models.TransformMarksheet.teacher_id == teacher.id
        ).order_by(models.TransformMarksheet.created_at.desc(), models.TransformMarksheet.id.desc()).all()
        for candidate in teacher_marksheets:
            course = db.query(models.Course).filter(models.Course.id == candidate.course_id).first()
            if course and normalize_text(course.course_name) == normalized_course_name:
                latest_marksheet = candidate
                break

    if latest_marksheet:
        auto_applied_updates, latest_payload = apply_reviewed_tasks_to_transform_marksheet(
            latest_marksheet,
            teacher,
            db
        )
        if auto_applied_updates > 0:
            auto_updated_sheet_id = latest_marksheet.id
            auto_updated_sheet_download_url = latest_payload.get("download_url", "")

    return {
        "status": "reviewed",
        "task": build_student_task_payload(task, teacher),
        "auto_updated_latest_marksheet": auto_applied_updates > 0,
        "auto_applied_updates": auto_applied_updates,
        "latest_marksheet_id": auto_updated_sheet_id,
        "latest_marksheet_download_url": auto_updated_sheet_download_url,
    }


@app.post("/api/admin/hods")
def create_hod(data: dict, db: Session = Depends(get_db)):
    full_name = normalize_text(data.get("full_name"))
    contact_no = normalize_text(data.get("contact_no"))
    department = normalize_text(data.get("department"), uppercase=True)

    if not all([full_name, department]):
        raise HTTPException(status_code=400, detail="Full name and department are required")

    existing = db.query(models.HeadOfDepartment).filter(
        or_(
            models.HeadOfDepartment.full_name == full_name,
            func.lower(models.HeadOfDepartment.department) == department.lower()
        )
    ).first()
    if existing:
        return {
            "status": "already_exists",
            "hod": {
                "id": existing.id,
                "full_name": existing.full_name,
                "username": existing.username,
                "password": existing.password,
                "hod_uid": existing.hod_uid,
                "department": existing.department
            }
        }

    username = generate_unique_username(full_name, models.HeadOfDepartment, "username", db)
    password = generate_easy_password(full_name, "hod")
    hod_uid = generate_unique_uid("HOD", models.HeadOfDepartment, "hod_uid", db)

    hod = models.HeadOfDepartment(
        hod_uid=hod_uid,
        username=username,
        full_name=full_name,
        password=password,
        contact_no=contact_no or None,
        department=department
    )
    db.add(hod)
    db.commit()
    db.refresh(hod)

    return {
        "status": "created",
        "hod": {
            "id": hod.id,
            "full_name": hod.full_name,
            "username": hod.username,
            "password": hod.password,
            "hod_uid": hod.hod_uid,
            "department": hod.department,
            "contact_no": hod.contact_no or ""
        }
    }


@app.post("/api/admin/students")
def create_student_from_admin(data: dict, db: Session = Depends(get_db)):
    full_name = normalize_text(data.get("full_name"))
    roll_no = normalize_text(data.get("roll_no"), uppercase=True)
    contact_no = normalize_text(data.get("contact_no"))
    department = normalize_text(data.get("department"), uppercase=True)
    program = normalize_text(data.get("program"))
    batch = normalize_text(data.get("batch"), uppercase=True)
    section = normalize_text(data.get("section"), uppercase=True)
    semester = data.get("semester")

    if not all([full_name, roll_no, department, program]):
        raise HTTPException(status_code=400, detail="Full name, roll number, department, and program are required")

    existing = db.query(models.Student).filter(
        or_(
            models.Student.roll_no == roll_no,
            models.Student.email == build_student_email(roll_no)
        )
    ).first()
    if existing:
        return {
            "status": "already_exists",
            "student": {
                "id": existing.id,
                "student_name": existing.student_name,
                "roll_no": existing.roll_no,
                "username": existing.username or "",
                "password": existing.password
            }
        }

    generated_email = build_student_email(roll_no)
    password = generate_easy_password(full_name, "std")

    student = models.Student(
        student_name=full_name,
        roll_no=roll_no,
        username=roll_no,
        password=password,
        contact_no=contact_no or None,
        email=generated_email,
        program=program or None,
        department=department or None,
        batch=batch or None,
        section=section or None,
        semester=int(semester) if semester not in (None, "") else None
    )
    db.add(student)
    db.commit()
    db.refresh(student)

    email_sent = send_credentials_email(
        generated_email,
        "TeachAssist Student Credentials",
        f"Student login ID: {student.roll_no}\nPassword: {student.password}\nEmail: {generated_email}"
    )

    return {
        "status": "created",
        "email_sent": email_sent,
        "student": {
            "id": student.id,
            "student_name": student.student_name,
            "roll_no": student.roll_no,
            "username": student.roll_no or "",
            "password": student.password,
            "email": student.email or "",
            "department": student.department or "",
            "program": student.program or ""
        }
      }


@app.post("/api/admin/students/import")
async def import_students_from_excel(file: UploadFile = File(...), db: Session = Depends(get_db)):
    filename = (file.filename or "").lower()
    if not filename.endswith(".xlsx"):
        raise HTTPException(status_code=400, detail="Please upload an Excel sheet in .xlsx format.")

    file_bytes = await file.read()
    if not file_bytes:
        raise HTTPException(status_code=400, detail="The uploaded Excel sheet is empty.")

    try:
        student_rows = parse_student_import_rows(file_bytes)
    except Exception as exc:
        raise HTTPException(status_code=400, detail=f"Could not read the Excel sheet: {exc}") from exc

    if not student_rows:
        raise HTTPException(status_code=400, detail="No valid student rows were found in the Excel sheet.")

    imported_count = 0
    updated_count = 0
    skipped_count = 0
    auxiliary_tables: set[str] = set()

    for row in student_rows:
        if not all([row["student_name"], row["roll_no"]]):
            skipped_count += 1
            continue

        generated_email = row["email"] or build_student_email(row["roll_no"])
        department = row["department"] or "GENERAL"
        batch = row["batch"] or "UNASSIGNED"
        program = row["program"] or department
        semester_value = row["semester"]
        semester = int(float(semester_value)) if semester_value not in ("", None) else None

        existing = db.query(models.Student).filter(models.Student.roll_no == row["roll_no"]).first()
        if existing:
            existing.student_name = row["student_name"]
            existing.username = row["roll_no"]
            existing.contact_no = row["contact_no"] or existing.contact_no
            existing.email = generated_email
            existing.program = program
            existing.department = department
            existing.batch = batch
            existing.section = row["section"] or existing.section
            existing.semester = semester
            student_record = existing
            updated_count += 1
        else:
            student_record = models.Student(
                student_name=row["student_name"],
                roll_no=row["roll_no"],
                username=row["roll_no"],
                password=generate_easy_password(row["student_name"], "std"),
                contact_no=row["contact_no"] or None,
                email=generated_email,
                program=program,
                department=department,
                batch=batch,
                section=row["section"] or None,
                semester=semester
            )
            db.add(student_record)
            db.flush()
            imported_count += 1

        table_name = build_student_batch_table_name(department, batch)
        ensure_student_batch_table(table_name, db)
        auxiliary_tables.add(table_name)
        mirror_student_to_batch_table(student_record, table_name, db)

    db.commit()

    return {
        "status": "imported",
        "imported_count": imported_count,
        "updated_count": updated_count,
        "skipped_count": skipped_count,
        "auxiliary_tables": sorted(auxiliary_tables),
    }


@app.post("/api/admin/courses")
def create_course_from_admin(data: dict, db: Session = Depends(get_db)):
    course_code = normalize_text(data.get("course_code"), uppercase=True)
    course_name = normalize_text(data.get("course_name"))

    if not all([course_code, course_name]):
        raise HTTPException(status_code=400, detail="Course code and course name are required")

    existing = db.query(models.Course).filter(
        func.lower(models.Course.course_code) == course_code.lower()
    ).first()
    if existing:
        return {
            "status": "already_exists",
            "course": {
                "id": existing.id,
                "course_code": existing.course_code,
                "course_name": existing.course_name
            }
        }

    course = models.Course(
        course_code=course_code,
        course_name=course_name
    )
    db.add(course)
    db.commit()
    db.refresh(course)

    return {
        "status": "created",
        "course": {
            "id": course.id,
            "course_code": course.course_code,
            "course_name": course.course_name
        }
    }


@app.post("/api/teacher/signup")
def teacher_signup(data: dict, db: Session = Depends(get_db)):
    full_name = normalize_text(data.get("full_name"))
    contact_no = normalize_text(data.get("contact_no"))
    email = normalize_text(data.get("email"))
    password = data.get("password") or ""

    if not full_name or not password.strip():
        raise HTTPException(status_code=400, detail="Full name and password are required")

    existing_teacher = db.query(models.Teacher).filter(
        or_(
            models.Teacher.email == email if email else False,
            models.Teacher.teacher_name == full_name
        )
    ).first()
    if existing_teacher:
        existing_teacher.contact_no = contact_no or existing_teacher.contact_no
        existing_teacher.email = email or existing_teacher.email
        existing_teacher.password = password.strip()
        existing_teacher.signup_source = existing_teacher.signup_source or "teacher_signup"
        if not existing_teacher.username:
            existing_teacher.username = generate_unique_username(full_name, models.Teacher, "username", db)
        if not existing_teacher.teacher_uid:
            existing_teacher.teacher_uid = generate_unique_uid("TCH", models.Teacher, "teacher_uid", db)
        db.commit()
        db.refresh(existing_teacher)
        return {
            "status": "updated",
            "teacher": {
                "teacher_name": existing_teacher.teacher_name,
                "username": existing_teacher.username,
                "password": "",
                "teacher_uid": existing_teacher.teacher_uid
            }
        }

    username = generate_unique_username(full_name, models.Teacher, "username", db)
    teacher_uid = generate_unique_uid("TCH", models.Teacher, "teacher_uid", db)

    teacher = models.Teacher(
        teacher_uid=teacher_uid,
        username=username,
        teacher_name=full_name,
        password=password.strip(),
        contact_no=contact_no or None,
        email=email or None,
        signup_source="teacher_signup"
    )
    db.add(teacher)

    signup_request = models.TeacherSignupRequest(
        full_name=full_name,
        roll_no=None,
        contact_no=contact_no or None,
        email=email or None,
        generated_username=username,
        generated_password="SELF_SET"
    )
    db.add(signup_request)
    db.commit()
    db.refresh(teacher)

    return {
        "status": "created",
        "teacher": {
            "teacher_name": teacher.teacher_name,
            "username": teacher.username,
            "password": "",
            "teacher_uid": teacher.teacher_uid
        }
    }


@app.post("/api/hod/teachers")
def create_teacher_from_hod(data: dict, db: Session = Depends(get_db)):
    hod_name = normalize_text(data.get("hod_name"))
    teacher_username = normalize_text(data.get("teacher_username"))
    department = normalize_text(data.get("department"), uppercase=True)
    course_code = normalize_text(data.get("course_code"), uppercase=True)

    if not all([hod_name, teacher_username, department, course_code]):
        raise HTTPException(status_code=400, detail="Teacher username, department, and course code are required")

    hod = db.query(models.HeadOfDepartment).filter(
        models.HeadOfDepartment.full_name == hod_name
    ).first()
    if not hod:
        raise HTTPException(status_code=404, detail="HOD not found")

    course = db.query(models.Course).filter(
        func.lower(models.Course.course_code) == course_code.lower()
    ).first()
    if not course:
        raise HTTPException(status_code=404, detail="Course not found")

    teacher = db.query(models.Teacher).filter(
        func.lower(models.Teacher.username) == teacher_username.lower()
    ).first()
    if not teacher:
        raise HTTPException(status_code=404, detail="Registered teacher not found for this username")

    if not teacher.department:
        teacher.department = department
        db.commit()
        db.refresh(teacher)

    existing_link = db.query(models.TeacherCourse).filter(
        models.TeacherCourse.teacher_id == teacher.id,
        models.TeacherCourse.course_id == course.id,
        models.TeacherCourse.department == department
    ).first()

    if existing_link:
        return {
            "status": "already_assigned",
            "teacher": {
                "id": teacher.id,
                "teacher_name": teacher.teacher_name,
                "teacher_uid": teacher.teacher_uid,
                "username": teacher.username,
                "password": "",
                "department": teacher.department or department,
                "course_code": course.course_code,
                "course_name": course.course_name
            }
        }

    db.add(models.TeacherCourse(
        teacher_id=teacher.id,
        course_id=course.id,
        department=department,
        threshold_percentage=50
    ))
    db.commit()

    return {
        "status": "assigned",
        "teacher": {
            "id": teacher.id,
            "teacher_name": teacher.teacher_name,
            "teacher_uid": teacher.teacher_uid,
            "username": teacher.username,
            "password": "",
            "department": teacher.department or department,
            "course_code": course.course_code,
            "course_name": course.course_name
        }
    }
# =========================================================
# REGISTER API
# =========================================================
@app.post("/api/register")
def register(data: dict, db: Session = Depends(get_db)):
    teacher_name = data.get("teacher_name")
    password = data.get("password")
    semester = data.get("semester")
    section = data.get("section", "")
    batch = data.get("batch", "")
    course_code = data.get("course_code")

    if not teacher_name or not password:
        raise HTTPException(status_code=400, detail="Missing fields")

    # Check if already exists
    existing = db.query(models.Teacher).filter(
        models.Teacher.teacher_name == teacher_name
    ).first()

    if existing:
        return {"status": "already_exists"}

    # Save teacher
    teacher = models.Teacher(
        teacher_uid=f"TCH-{uuid4().hex[:8].upper()}",
        username=teacher_name.lower().replace(" ", "."),
        teacher_name=teacher_name,
        password=password
    )
    db.add(teacher)
    db.commit()
    db.refresh(teacher)

    # Save course mapping if course_code given
    if course_code and semester:
        course = db.query(models.Course).filter(
            models.Course.course_code == course_code
        ).first()

        if course:
            tc = models.TeacherCourse(
                teacher_id=teacher.id,
                course_id=course.id,
                semester=int(semester),
                batch=batch
            )
            db.add(tc)
            db.commit()

    return {"status": "registered", "teacher_name": teacher_name}
# =========================================================
# COURSE LOOKUP BY CODE
# =========================================================
@app.get("/api/courses/lookup")
def lookup_course(course_code: str, db: Session = Depends(get_db)):
    course = db.query(models.Course).filter(
        func.lower(models.Course.course_code) == course_code.lower().strip()
    ).first()
    
    if not course:
        raise HTTPException(status_code=404, detail="Course not found")
    
    return {
        "course_code": course.course_code,
        "course_name": course.course_name
    }


@app.get("/api/courses/search")
def search_courses(q: str = "", db: Session = Depends(get_db)):
    query = normalize_text(q, uppercase=True)
    courses_query = db.query(models.Course)

    if query:
        like_value = f"%{query}%"
        courses_query = courses_query.filter(
            or_(
                func.upper(models.Course.course_code).like(like_value),
                func.upper(models.Course.course_name).like(like_value)
            )
        )

    courses = courses_query.order_by(models.Course.course_code.asc()).limit(12).all()
    return {
        "courses": [
            {
                "course_code": course.course_code,
                "course_name": course.course_name
            }
            for course in courses
        ]
    }

# =========================================================
# TEACHER SETUP
# =========================================================
@app.post("/api/teacher/setup")
def teacher_setup(data: dict, db: Session = Depends(get_db)):
    return create_teacher_record(data, db)


@app.post("/api/teacher/profile/records")
def add_teacher_profile_record(data: dict, db: Session = Depends(get_db)):
    return create_teacher_record(data, db)


@app.post("/api/teacher/profile/records/{record_id}/threshold")
def edit_teacher_profile_threshold(record_id: int, data: dict, db: Session = Depends(get_db)):
    return update_teacher_record_threshold(record_id, data, db)


@app.post("/api/teacher/profile/records/{record_id}")
def edit_teacher_profile_record(record_id: int, data: dict, db: Session = Depends(get_db)):
    return update_teacher_record(record_id, data, db)


@app.delete("/api/teacher/profile/records/{record_id}")
def remove_teacher_profile_record(record_id: int, teacher_name: str, db: Session = Depends(get_db)):
    return delete_teacher_record(record_id, teacher_name, db)


@app.get("/api/teacher/profile")
def teacher_profile(teacher_name: str, db: Session = Depends(get_db)):
    teacher = db.query(models.Teacher).filter(
        models.Teacher.teacher_name == teacher_name
    ).first()

    if not teacher:
        raise HTTPException(status_code=404, detail="Teacher not found")

    return get_teacher_profile_data(teacher, db)



def get_transform_excel_builder():
    builder = load_t2_excel_builder()
    if builder is None:
        raise HTTPException(status_code=500, detail="Transform Excel builder is not available")
    return builder


def normalize_clo_label(value: str) -> str:
    return re.sub(r"[^A-Z0-9]+", "", (value or "").upper())


def distribute_score_across_labels(
    row_marks: dict[str, float],
    matching_labels: list[str],
    assessment_totals: dict,
    score_value: float
) -> None:
    remaining = max(0, int(round(score_value)))
    eligible = []
    for label in matching_labels:
        max_marks = int(round(float(assessment_totals.get(label, row_marks.get(label, 0) or 0) or 0)))
        current_value = int(round(float(row_marks.get(label, 0) or 0)))
        room = max(0, max_marks - current_value)
        if room > 0:
            eligible.append(label)

    while remaining > 0 and eligible:
        share = remaining // len(eligible)
        remainder = remaining % len(eligible)
        progressed = False
        next_eligible: list[str] = []

        for index, label in enumerate(eligible):
            max_marks = int(round(float(assessment_totals.get(label, row_marks.get(label, 0) or 0) or 0)))
            current_value = int(round(float(row_marks.get(label, 0) or 0)))
            room = max(0, max_marks - current_value)
            target_add = share + (1 if index < remainder else 0)
            actual_add = min(room, target_add)
            if actual_add > 0:
                progressed = True
                row_marks[label] = current_value + actual_add
                remaining -= actual_add
            else:
                row_marks[label] = current_value

            room_after = max(0, max_marks - int(round(float(row_marks.get(label, 0) or 0))))
            if room_after > 0:
                next_eligible.append(label)

        if not progressed:
            break
        eligible = next_eligible


def apply_reviewed_tasks_to_transform_marksheet(
    marksheet: models.TransformMarksheet,
    teacher: models.Teacher,
    db: Session
) -> tuple[int, dict]:
    base_payload = build_transform_marksheet_payload(marksheet, db)
    selected_options = base_payload["selected_options"]
    assessment_totals = base_payload["assessment_totals"]
    marksheet_course_code = normalize_text(base_payload.get("course_code") or "", uppercase=True)
    marksheet_course_name = normalize_text(base_payload.get("course_name") or "")

    student_rows_by_id = {
        row["student_id"]: {
            "student_id": row["student_id"],
            "roll_no": row["roll_no"],
            "student_name": row["student_name"],
            "marks": dict(row["obtained_marks_by_assessment"]),
            "remarks": row.get("remarks") or ""
        }
        for row in base_payload["student_marks"]
    }
    student_rows_by_roll = {
        normalize_text(row["roll_no"], uppercase=True): row
        for row in student_rows_by_id.values()
        if normalize_text(row.get("roll_no"), uppercase=True)
    }

    reviewed_tasks = db.query(models.StudentTask).filter(
        models.StudentTask.teacher_id == teacher.id,
        models.StudentTask.reviewed_at.isnot(None)
    ).order_by(models.StudentTask.reviewed_at.asc(), models.StudentTask.id.asc()).all()

    applied_updates = 0
    for task in reviewed_tasks:
        task_course_matches = False
        if task.course_id and task.course_id == marksheet.course_id:
            task_course_matches = True
        elif not task.course_id:
            task_course_name = normalize_text(task.course_name_snapshot or "")
            task_title = normalize_text(task.title or "")
            if task_course_name and marksheet_course_name and task_course_name == marksheet_course_name:
                task_course_matches = True
            elif marksheet_course_code and task_title and marksheet_course_code in task_title.upper():
                task_course_matches = True

        if not task_course_matches:
            continue

        task_roll_no = normalize_text(task.assigned_roll_no, uppercase=True)
        if not task_roll_no:
            task_student = db.query(models.Student).filter(models.Student.id == task.student_id).first()
            task_roll_no = normalize_text(task_student.roll_no if task_student else "", uppercase=True)

        row = student_rows_by_roll.get(task_roll_no)
        if not row:
            continue
        clo_key = normalize_clo_label(task.clo or "")
        if not clo_key:
            continue
        task_token = f"[task:{task.id}]"
        existing_remarks = row.get("remarks", "") or ""
        if task_token in existing_remarks:
            continue

        score_value = None
        try:
            score_value = float(task.teacher_score) if task.teacher_score not in (None, "") else None
        except (TypeError, ValueError):
            score_value = None

        matching_labels = [
            label for label in row["marks"].keys()
            if clo_key and clo_key in normalize_clo_label(label)
        ]
        if not matching_labels:
            continue

        if score_value is not None:
            distribute_score_across_labels(row["marks"], matching_labels, assessment_totals, score_value)

        status_note = f"{task.clo or 'Task'} {task.teacher_decision or ''} {task_token}".strip()
        row["remarks"] = "; ".join(part for part in [existing_remarks.strip(), status_note] if part)
        applied_updates += 1

    if applied_updates == 0:
        return 0, base_payload

    db.query(models.TransformStudentAssessmentMark).filter(
        models.TransformStudentAssessmentMark.marksheet_id == marksheet.id
    ).delete(synchronize_session=False)
    db.commit()

    excel_rows = []
    for row in student_rows_by_id.values():
        for label, value in row["marks"].items():
            db.add(models.TransformStudentAssessmentMark(
                marksheet_id=marksheet.id,
                student_id=row["student_id"],
                assessment_label=label,
                obtained_marks=int(round(float(value or 0))),
                remarks=(row.get("remarks") or "").strip() or None
            ))
        excel_rows.append({
            "roll_number": row["roll_no"],
            "full_name": row["student_name"],
            "marks": row["marks"],
            "remarks": row.get("remarks") or "",
        })
    db.commit()

    file_path = write_transform_marksheet_excel(
        marksheet.id,
        teacher.teacher_name,
        teacher.email or "",
        marksheet.department or teacher.department or "",
        marksheet.section or "",
        marksheet.total_marks,
        selected_options,
        assessment_totals,
        excel_rows,
    )
    marksheet.excel_file_path = file_path
    db.commit()
    db.refresh(marksheet)

    return applied_updates, build_transform_marksheet_payload(marksheet, db)


def parse_transform_json(value: str | None, fallback: dict | None = None) -> dict:
    if not value:
        return fallback or {}
    try:
        parsed = json.loads(value)
        return parsed if isinstance(parsed, dict) else (fallback or {})
    except Exception:
        return fallback or {}


def get_transform_marksheet_marks(marksheet_id: int, db: Session) -> list[models.TransformStudentAssessmentMark]:
    return db.query(models.TransformStudentAssessmentMark).filter(
        models.TransformStudentAssessmentMark.marksheet_id == marksheet_id
    ).all()


def transform_marksheet_expires_at(marksheet: models.TransformMarksheet) -> datetime | None:
    if not marksheet.created_at:
        return None
    created_at = marksheet.created_at
    if created_at.tzinfo is None:
        created_at = created_at.replace(tzinfo=timezone.utc)
    return created_at + timedelta(days=TRANSFORM_MARKSHEET_RETENTION_DAYS)


def transform_marksheet_is_expired(marksheet: models.TransformMarksheet) -> bool:
    expires_at = transform_marksheet_expires_at(marksheet)
    return bool(expires_at and datetime.now(timezone.utc) >= expires_at)


def delete_transform_marksheet_record(marksheet: models.TransformMarksheet, db: Session) -> None:
    db.query(models.TransformMarksheet).filter(
        models.TransformMarksheet.source_marksheet_id == marksheet.id
    ).update(
        {models.TransformMarksheet.source_marksheet_id: None},
        synchronize_session=False
    )

    db.query(models.TransformStudentAssessmentMark).filter(
        models.TransformStudentAssessmentMark.marksheet_id == marksheet.id
    ).delete(synchronize_session=False)

    excel_path = marksheet.excel_file_path
    db.delete(marksheet)
    db.commit()

    if excel_path:
        try:
            path = Path(excel_path)
            if path.exists():
                path.unlink()
        except Exception:
            pass


def cleanup_expired_transform_marksheets(db: Session, teacher_id: int | None = None) -> None:
    cutoff = datetime.now(timezone.utc) - timedelta(days=TRANSFORM_MARKSHEET_RETENTION_DAYS)
    query = db.query(models.TransformMarksheet).filter(
        models.TransformMarksheet.created_at.isnot(None),
        models.TransformMarksheet.created_at < cutoff
    )
    if teacher_id is not None:
        query = query.filter(models.TransformMarksheet.teacher_id == teacher_id)

    for marksheet in query.all():
        delete_transform_marksheet_record(marksheet, db)


def build_transform_marksheet_payload(marksheet: models.TransformMarksheet, db: Session) -> dict:
    course = db.query(models.Course).filter(models.Course.id == marksheet.course_id).first()
    record = db.query(models.TeacherCourse).filter(models.TeacherCourse.id == marksheet.teacher_course_id).first()
    marks = get_transform_marksheet_marks(marksheet.id, db)
    student_ids = list({mark.student_id for mark in marks})
    student_map = {
        student.id: student
        for student in db.query(models.Student).filter(models.Student.id.in_(student_ids)).all()
    } if student_ids else {}

    selected_options = parse_transform_json(marksheet.selected_options, {})
    assessment_totals = parse_transform_json(marksheet.assessment_totals, {})

    grouped_rows: dict[int, dict] = {}
    for mark in marks:
        student = student_map.get(mark.student_id)
        if not student:
            continue
        row = grouped_rows.setdefault(mark.student_id, {
            "student_id": student.id,
            "student_name": student.student_name,
            "roll_no": student.roll_no or "",
            "obtained_marks_by_assessment": {},
            "remarks": mark.remarks or ""
        })
        row["obtained_marks_by_assessment"][mark.assessment_label] = float(mark.obtained_marks)
        if mark.remarks:
            row["remarks"] = mark.remarks

    return {
        "id": marksheet.id,
        "teacher_course_id": marksheet.teacher_course_id,
        "course_id": marksheet.course_id,
        "course_code": course.course_code if course else "",
        "course_name": course.course_name if course else "",
        "semester": marksheet.semester,
        "section": marksheet.section or "",
        "batch": marksheet.batch or "",
        "department": marksheet.department or "",
        "exam_type": marksheet.exam_type or "Midterm",
        "total_marks": marksheet.total_marks,
        "export_file_name": marksheet.export_file_name or "",
        "selected_options": selected_options,
        "assessment_totals": assessment_totals,
        "download_url": f"/api/transform/marksheets/{marksheet.id}/download",
        "created_at": marksheet.created_at.isoformat() if marksheet.created_at else None,
        "expires_at": transform_marksheet_expires_at(marksheet).isoformat() if transform_marksheet_expires_at(marksheet) else None,
        "source_kind": marksheet.source_kind,
        "source_marksheet_id": marksheet.source_marksheet_id,
        "student_count": len(grouped_rows),
        "student_marks": list(grouped_rows.values()),
        "teacher_threshold_percentage": record.threshold_percentage if record and record.threshold_percentage is not None else 50
    }


def build_transform_marksheet_summary_payload(marksheet: models.TransformMarksheet, db: Session) -> dict:
    course = db.query(models.Course).filter(models.Course.id == marksheet.course_id).first()
    record = db.query(models.TeacherCourse).filter(models.TeacherCourse.id == marksheet.teacher_course_id).first()
    student_count = db.query(
        func.count(func.distinct(models.TransformStudentAssessmentMark.student_id))
    ).filter(
        models.TransformStudentAssessmentMark.marksheet_id == marksheet.id
    ).scalar() or 0

    return {
        "id": marksheet.id,
        "teacher_course_id": marksheet.teacher_course_id,
        "course_id": marksheet.course_id,
        "course_code": course.course_code if course else "",
        "course_name": course.course_name if course else "",
        "semester": marksheet.semester,
        "section": marksheet.section or "",
        "batch": marksheet.batch or "",
        "department": marksheet.department or "",
        "exam_type": marksheet.exam_type or "Midterm",
        "total_marks": marksheet.total_marks,
        "export_file_name": marksheet.export_file_name or "",
        "selected_options": parse_transform_json(marksheet.selected_options, {}),
        "assessment_totals": parse_transform_json(marksheet.assessment_totals, {}),
        "download_url": f"/api/transform/marksheets/{marksheet.id}/download",
        "created_at": marksheet.created_at.isoformat() if marksheet.created_at else None,
        "expires_at": transform_marksheet_expires_at(marksheet).isoformat() if transform_marksheet_expires_at(marksheet) else None,
        "source_kind": marksheet.source_kind,
        "source_marksheet_id": marksheet.source_marksheet_id,
        "student_count": student_count,
        "teacher_threshold_percentage": record.threshold_percentage if record and record.threshold_percentage is not None else 50
    }


def write_transform_marksheet_excel(
    marksheet_id: int,
    teacher_name: str,
    teacher_email: str,
    department_name: str,
    section: str,
    total_marks: int,
    selected_options: dict,
    assessment_totals: dict,
    students: list[dict]
) -> str:
    builder = get_transform_excel_builder()
    assessment_columns = list(selected_options.keys())
    return builder(
        sheet_id=marksheet_id,
        teacher_name=teacher_name,
        teacher_email=teacher_email,
        department=department_name,
        section=section,
        total_marks=total_marks,
        selected_options=selected_options,
        assessment_totals=assessment_totals,
        assessment_columns=assessment_columns,
        students=students,
    )


def persist_transform_marksheet_rows(
    marksheet: models.TransformMarksheet,
    teacher: models.Teacher,
    record: models.TeacherCourse,
    selected_options: dict,
    assessment_totals: dict,
    student_marks: list[dict],
    db: Session
) -> None:
    students = get_students_for_teacher_record(record, db)
    student_map = {student.id: student for student in students}

    db.query(models.TransformStudentAssessmentMark).filter(
        models.TransformStudentAssessmentMark.marksheet_id == marksheet.id
    ).delete(synchronize_session=False)
    db.commit()

    excel_rows = []
    for row in student_marks:
        student_id = row.get("student_id")
        if student_id not in student_map:
            continue

        marks_map = {}
        for mark in row.get("marks", []):
            label = str(mark.get("assessment_label") or "").strip()
            if not label:
                continue
            obtained_marks = float(mark.get("obtained_marks") or 0)
            db.add(models.TransformStudentAssessmentMark(
                marksheet_id=marksheet.id,
                student_id=student_id,
                assessment_label=label,
                obtained_marks=obtained_marks,
                remarks=(row.get("remarks") or "").strip() or None
            ))
            marks_map[label] = obtained_marks

        if not marks_map and not (row.get("remarks") or "").strip():
            continue

        student = student_map[student_id]
        excel_rows.append({
            "roll_number": student.roll_no or "",
            "full_name": student.student_name,
            "marks": marks_map,
            "remarks": (row.get("remarks") or "").strip(),
        })

    db.commit()

    file_path = write_transform_marksheet_excel(
        marksheet.id,
        teacher.teacher_name,
        teacher.email or "",
        record.department or teacher.department or "",
        record.section or "",
        marksheet.total_marks,
        selected_options,
        assessment_totals,
        excel_rows,
    )
    marksheet.excel_file_path = file_path
    db.commit()
    db.refresh(marksheet)


def get_teacher_course_transform_options(teacher: models.Teacher, db: Session) -> list[dict]:
    records = db.query(models.TeacherCourse).filter(
        models.TeacherCourse.teacher_id == teacher.id
    ).order_by(models.TeacherCourse.created_at.desc(), models.TeacherCourse.id.desc()).all()

    course_ids = [record.course_id for record in records if record.course_id]
    course_map = {
        course.id: course
        for course in db.query(models.Course).filter(models.Course.id.in_(course_ids)).all()
    } if course_ids else {}

    options = []
    for record in records:
        students = get_students_for_teacher_record(record, db)
        latest_marksheet = db.query(models.TransformMarksheet).filter(
            models.TransformMarksheet.teacher_course_id == record.id
        ).order_by(models.TransformMarksheet.created_at.desc(), models.TransformMarksheet.id.desc()).first()
        reviewed_task_count = db.query(models.StudentTask).filter(
            models.StudentTask.teacher_id == teacher.id,
            models.StudentTask.course_id == record.course_id,
            models.StudentTask.reviewed_at.isnot(None)
        ).count()
        course = course_map.get(record.course_id)
        options.append({
            "record_id": record.id,
            "course_id": record.course_id,
            "course_code": course.course_code if course else "",
            "course_name": course.course_name if course else "",
            "semester": record.semester,
            "section": record.section or "",
            "batch": record.batch or "",
            "department": record.department or "",
            "student_count": len(students),
            "threshold_percentage": record.threshold_percentage if record.threshold_percentage is not None else 50,
            "latest_marksheet_id": latest_marksheet.id if latest_marksheet else None,
            "reviewed_task_count": reviewed_task_count,
        })
    return options


@app.get("/api/transform/courses")
def get_transform_courses(teacher_name: str, db: Session = Depends(get_db)):
    teacher = db.query(models.Teacher).filter(
        models.Teacher.teacher_name == teacher_name
    ).first()
    if not teacher:
        raise HTTPException(status_code=404, detail="Teacher not found")

    return {
        "teacher_name": teacher.teacher_name,
        "teacher_email": teacher.email or "",
        "teacher_uid": teacher.teacher_uid or "",
        "department": teacher.department or "",
        "courses": get_teacher_course_transform_options(teacher, db)
    }


@app.get("/api/transform/courses/{record_id}/students")
def get_transform_course_students(record_id: int, teacher_name: str, db: Session = Depends(get_db)):
    teacher = db.query(models.Teacher).filter(models.Teacher.teacher_name == teacher_name).first()
    if not teacher:
        raise HTTPException(status_code=404, detail="Teacher not found")

    record = db.query(models.TeacherCourse).filter(
        models.TeacherCourse.id == record_id,
        models.TeacherCourse.teacher_id == teacher.id
    ).first()
    if not record:
        raise HTTPException(status_code=404, detail="Course record not found")

    course = db.query(models.Course).filter(models.Course.id == record.course_id).first()
    students = get_students_for_teacher_record(record, db)
    latest_marksheet = db.query(models.TransformMarksheet).filter(
        models.TransformMarksheet.teacher_course_id == record.id
    ).order_by(models.TransformMarksheet.created_at.desc(), models.TransformMarksheet.id.desc()).first()

    return {
        "teacher_name": teacher.teacher_name,
        "teacher_email": teacher.email or "",
        "teacher_uid": teacher.teacher_uid or "",
        "course": {
            "record_id": record.id,
            "course_id": record.course_id,
            "course_code": course.course_code if course else "",
            "course_name": course.course_name if course else "",
            "semester": record.semester,
            "section": record.section or "",
            "batch": record.batch or "",
            "department": record.department or "",
            "student_count": len(students),
            "threshold_percentage": record.threshold_percentage if record.threshold_percentage is not None else 50,
            "latest_marksheet_id": latest_marksheet.id if latest_marksheet else None,
        },
        "students": [
            {
                "id": student.id,
                "roll_number": student.roll_no or "",
                "full_name": student.student_name,
                "batch": student.batch or "",
                "section": student.section or "",
                "semester": str(student.semester) if student.semester is not None else None,
            }
            for student in students
        ]
    }


@app.post("/api/transform/marksheets")
def save_transform_marksheet(data: dict, db: Session = Depends(get_db)):
    teacher_name = normalize_text(data.get("teacher_name"))
    teacher_course_id = data.get("teacher_course_id")
    total_marks = int(float(data.get("total_marks") or 100))
    export_file_name = (data.get("export_file_name") or "").strip() or "marksheet"
    selected_options = data.get("selected_options") or {}
    assessment_totals = data.get("assessment_totals") or {}
    student_marks = data.get("student_marks") or []
    exam_type = (data.get("exam_type") or "Midterm").strip()
    source_kind = (data.get("source_kind") or "initial").strip() or "initial"
    source_marksheet_id = data.get("source_marksheet_id")

    teacher = db.query(models.Teacher).filter(models.Teacher.teacher_name == teacher_name).first()
    if not teacher:
        raise HTTPException(status_code=404, detail="Teacher not found")
    cleanup_expired_transform_marksheets(db, teacher.id)

    record = db.query(models.TeacherCourse).filter(
        models.TeacherCourse.id == teacher_course_id,
        models.TeacherCourse.teacher_id == teacher.id
    ).first()
    if not record:
        raise HTTPException(status_code=404, detail="Teacher course record not found")

    course = db.query(models.Course).filter(models.Course.id == record.course_id).first()

    if not selected_options:
        raise HTTPException(status_code=400, detail="Please add at least one assessment column")
    marksheet = models.TransformMarksheet(
        teacher_id=teacher.id,
        teacher_course_id=record.id,
        course_id=record.course_id,
        semester=record.semester,
        section=record.section,
        batch=record.batch,
        department=record.department,
        exam_type=exam_type,
        total_marks=total_marks,
        export_file_name=export_file_name,
        selected_options=json.dumps(selected_options),
        assessment_totals=json.dumps(assessment_totals),
        source_kind=source_kind,
        source_marksheet_id=source_marksheet_id,
    )
    db.add(marksheet)
    db.commit()
    db.refresh(marksheet)

    persist_transform_marksheet_rows(
        marksheet=marksheet,
        teacher=teacher,
        record=record,
        selected_options=selected_options,
        assessment_totals=assessment_totals,
        student_marks=student_marks,
        db=db,
    )

    payload = build_transform_marksheet_payload(marksheet, db)
    return {
        "message": "Marksheet saved successfully.",
        "marksheet": payload,
        "download_url": payload["download_url"],
        "course_code": course.course_code if course else "",
        "course_name": course.course_name if course else ""
    }


@app.put("/api/transform/marksheets/{marksheet_id}")
def update_transform_marksheet(marksheet_id: int, data: dict, db: Session = Depends(get_db)):
    teacher_name = normalize_text(data.get("teacher_name"))
    teacher_course_id = data.get("teacher_course_id")
    total_marks = int(float(data.get("total_marks") or 100))
    export_file_name = (data.get("export_file_name") or "").strip() or "marksheet"
    selected_options = data.get("selected_options") or {}
    assessment_totals = data.get("assessment_totals") or {}
    student_marks = data.get("student_marks") or []
    exam_type = (data.get("exam_type") or "Midterm").strip()
    source_kind = (data.get("source_kind") or "initial").strip() or "initial"
    source_marksheet_id = data.get("source_marksheet_id")

    teacher = db.query(models.Teacher).filter(models.Teacher.teacher_name == teacher_name).first()
    if not teacher:
        raise HTTPException(status_code=404, detail="Teacher not found")
    cleanup_expired_transform_marksheets(db, teacher.id)

    marksheet = db.query(models.TransformMarksheet).filter(
        models.TransformMarksheet.id == marksheet_id,
        models.TransformMarksheet.teacher_id == teacher.id
    ).first()
    if not marksheet:
        raise HTTPException(status_code=404, detail="Marksheet not found")
    if transform_marksheet_is_expired(marksheet):
        delete_transform_marksheet_record(marksheet, db)
        raise HTTPException(status_code=404, detail="Marksheet expired after 28 days and was deleted")

    record = db.query(models.TeacherCourse).filter(
        models.TeacherCourse.id == teacher_course_id,
        models.TeacherCourse.teacher_id == teacher.id
    ).first()
    if not record:
        raise HTTPException(status_code=404, detail="Teacher course record not found")

    course = db.query(models.Course).filter(models.Course.id == record.course_id).first()
    if not selected_options:
        raise HTTPException(status_code=400, detail="Please add at least one assessment column")
    marksheet.teacher_course_id = record.id
    marksheet.course_id = record.course_id
    marksheet.semester = record.semester
    marksheet.section = record.section
    marksheet.batch = record.batch
    marksheet.department = record.department
    marksheet.exam_type = exam_type
    marksheet.total_marks = total_marks
    marksheet.export_file_name = export_file_name
    marksheet.selected_options = json.dumps(selected_options)
    marksheet.assessment_totals = json.dumps(assessment_totals)
    marksheet.source_kind = source_kind
    marksheet.source_marksheet_id = source_marksheet_id
    db.commit()
    db.refresh(marksheet)

    persist_transform_marksheet_rows(
        marksheet=marksheet,
        teacher=teacher,
        record=record,
        selected_options=selected_options,
        assessment_totals=assessment_totals,
        student_marks=student_marks,
        db=db,
    )

    payload = build_transform_marksheet_payload(marksheet, db)
    return {
        "message": "Marksheet updated successfully.",
        "marksheet": payload,
        "download_url": payload["download_url"],
        "course_code": course.course_code if course else "",
        "course_name": course.course_name if course else ""
    }


@app.get("/api/transform/marksheets")
def list_transform_marksheets(teacher_name: str, db: Session = Depends(get_db)):
    teacher = db.query(models.Teacher).filter(models.Teacher.teacher_name == teacher_name).first()
    if not teacher:
        raise HTTPException(status_code=404, detail="Teacher not found")
    cleanup_expired_transform_marksheets(db, teacher.id)

    marksheets = db.query(models.TransformMarksheet).filter(
        models.TransformMarksheet.teacher_id == teacher.id
    ).order_by(
        models.TransformMarksheet.created_at.desc(),
        models.TransformMarksheet.id.desc()
    ).all()

    return {
        "marksheets": [build_transform_marksheet_summary_payload(sheet, db) for sheet in marksheets]
    }


@app.get("/api/transform/marksheets/{marksheet_id}")
def get_transform_marksheet_detail(marksheet_id: int, teacher_name: str, db: Session = Depends(get_db)):
    teacher = db.query(models.Teacher).filter(models.Teacher.teacher_name == teacher_name).first()
    if not teacher:
        raise HTTPException(status_code=404, detail="Teacher not found")
    cleanup_expired_transform_marksheets(db, teacher.id)

    marksheet = db.query(models.TransformMarksheet).filter(
        models.TransformMarksheet.id == marksheet_id,
        models.TransformMarksheet.teacher_id == teacher.id
    ).first()
    if not marksheet:
        raise HTTPException(status_code=404, detail="Marksheet not found")
    if transform_marksheet_is_expired(marksheet):
        delete_transform_marksheet_record(marksheet, db)
        raise HTTPException(status_code=404, detail="Marksheet expired after 28 days and was deleted")

    return build_transform_marksheet_payload(marksheet, db)


@app.get("/api/transform/marksheets/{marksheet_id}/download")
def download_transform_marksheet(marksheet_id: int, db: Session = Depends(get_db)):
    marksheet = db.query(models.TransformMarksheet).filter(
        models.TransformMarksheet.id == marksheet_id
    ).first()
    if not marksheet or not marksheet.excel_file_path:
        raise HTTPException(status_code=404, detail="Marksheet file not found")
    if transform_marksheet_is_expired(marksheet):
        delete_transform_marksheet_record(marksheet, db)
        raise HTTPException(status_code=404, detail="Marksheet expired after 28 days and was deleted")
    return FileResponse(
        path=marksheet.excel_file_path,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        filename=f"{marksheet.export_file_name or f'transform_marksheet_{marksheet.id}'}.xlsx"
    )


@app.delete("/api/transform/marksheets/{marksheet_id}")
def delete_transform_marksheet(marksheet_id: int, teacher_name: str, db: Session = Depends(get_db)):
    teacher = db.query(models.Teacher).filter(models.Teacher.teacher_name == teacher_name).first()
    if not teacher:
        raise HTTPException(status_code=404, detail="Teacher not found")
    cleanup_expired_transform_marksheets(db, teacher.id)

    marksheet = db.query(models.TransformMarksheet).filter(
        models.TransformMarksheet.id == marksheet_id,
        models.TransformMarksheet.teacher_id == teacher.id
    ).first()
    if not marksheet:
        raise HTTPException(status_code=404, detail="Marksheet not found")
    if transform_marksheet_is_expired(marksheet):
        delete_transform_marksheet_record(marksheet, db)
        raise HTTPException(status_code=404, detail="Marksheet expired after 28 days and was deleted")

    delete_transform_marksheet_record(marksheet, db)

    return {"status": "deleted", "marksheet_id": marksheet_id}


@app.post("/api/transform/marksheets/{marksheet_id}/updated")
def generate_updated_transform_marksheet(marksheet_id: int, data: dict, db: Session = Depends(get_db)):
    teacher_name = normalize_text(data.get("teacher_name"))
    teacher = db.query(models.Teacher).filter(models.Teacher.teacher_name == teacher_name).first()
    if not teacher:
        raise HTTPException(status_code=404, detail="Teacher not found")
    cleanup_expired_transform_marksheets(db, teacher.id)

    base_marksheet = db.query(models.TransformMarksheet).filter(
        models.TransformMarksheet.id == marksheet_id,
        models.TransformMarksheet.teacher_id == teacher.id
    ).first()
    if not base_marksheet:
        raise HTTPException(status_code=404, detail="Base marksheet not found")
    if transform_marksheet_is_expired(base_marksheet):
        delete_transform_marksheet_record(base_marksheet, db)
        raise HTTPException(status_code=404, detail="Marksheet expired after 28 days and was deleted")

    applied_updates, updated_payload = apply_reviewed_tasks_to_transform_marksheet(base_marksheet, teacher, db)
    if applied_updates == 0:
        raise HTTPException(
            status_code=400,
            detail="No new reviewed task changes were available for this sheet."
        )
    return {
        "message": "Latest marksheet updated successfully.",
        "applied_updates": applied_updates,
        "marksheet": updated_payload,
        "download_url": updated_payload["download_url"],
        "course_code": updated_payload["course_code"],
        "course_name": updated_payload["course_name"]
    }

# =========================================================
# CHECK TEACHER SETUP STATUS
# =========================================================
@app.get("/api/teacher/status")
def teacher_status(teacher_name: str, db: Session = Depends(get_db)):
    teacher = db.query(models.Teacher).filter(
        models.Teacher.teacher_name == teacher_name
    ).first()

    if not teacher:
        return {"status": "not_found"}

    tc_count = db.query(models.TeacherCourse).filter(
        models.TeacherCourse.teacher_id == teacher.id
    ).count()

    if tc_count > 0:
        return {"status": "setup_complete"}
    else:
        return {"status": "needs_setup"}

# ------------------ LLM + RAG ------------------
rag = RAGEngine()
groq_client = Groq(api_key=os.getenv("GROQ_API_KEY"))
GROQ_MODEL = "llama-3.1-8b-instant"

# =========================================================
# FILE TEXT EXTRACTION
# =========================================================
def extract_text_from_file(upload_file: UploadFile) -> str:
    raw = upload_file.file.read()
    filename = upload_file.filename.lower()

    if filename.endswith(".pdf"):
        reader = PdfReader(BytesIO(raw))
        return "\n".join(page.extract_text() or "" for page in reader.pages)

    if filename.endswith(".pptx"):
        prs = Presentation(BytesIO(raw))
        slides_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slides_text.append(shape.text)
        return "\n".join(slides_text)

    if filename.endswith(".docx"):
        doc = Document(BytesIO(raw))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

    if filename.endswith(".txt"):
        return raw.decode("utf-8", errors="ignore")

    raise HTTPException(
        status_code=400,
        detail="Unsupported file format. Upload PDF, PPTX, DOCX, or TXT."
    )

# ------------------ CLEAN OUTPUT ------------------
def clean_output(text: str) -> str:
    text = re.sub(r"\*{1,3}", "", text)
    text = text.replace("•", "-").replace("###", "")
    return text.strip()

# ------------------ TEXT CHUNKING ------------------
def chunk_text(text: str, size=800, overlap=100):
    chunks, start = [], 0
    while start < len(text):
        chunks.append(text[start:start + size])
        start += size - overlap
    return chunks


def clip_text(text: str, max_chars: int) -> str:
    if len(text) <= max_chars:
        return text
    return text[:max_chars].rstrip() + "\n...[truncated]"

# ------------------ EXTRACT MARKS FROM PROMPT ------------------
def extract_marks_config(prompt: str) -> Dict[str, int]:
    """Extract marks configuration from the prompt"""
    marks = {
        'mcq': 1,
        'short': 5,
        'long': 10,
        'scenario': 10,
        'task': 4  # Default for assignment tasks
    }
    
    # Try to extract marks from prompt
    mcq_match = re.search(r'MCQ.*?(\d+)\s*marks', prompt, re.IGNORECASE)
    short_match = re.search(r'Short.*?(\d+)\s*marks', prompt, re.IGNORECASE)
    long_match = re.search(r'Long.*?(\d+)\s*marks', prompt, re.IGNORECASE)
    scenario_match = re.search(r'Scenario.*?(\d+)\s*marks', prompt, re.IGNORECASE)
    task_match = re.search(r'MARKS_PER_TASK:\s*(\d+)', prompt, re.IGNORECASE)
    
    if mcq_match:
        marks['mcq'] = int(mcq_match.group(1))
    if short_match:
        marks['short'] = int(short_match.group(1))
    if long_match:
        marks['long'] = int(long_match.group(1))
    if scenario_match:
        marks['scenario'] = int(scenario_match.group(1))
    if task_match:
        marks['task'] = int(task_match.group(1))
    
    return marks

# ------------------ PARSE EXAM STRUCTURE ------------------
def parse_exam_for_export(content: str, marks_config: Dict[str, int], exam_type: str) -> Dict[str, Any]:
    """Parse exam content for structured export"""
    lines = content.split('\n')
    sections = []
    current_section = None
    current_qtype = None
    answer_key_started = False
    
    questions_part = []
    answers_part = []
    tasks = []  # For assignments
    
    for line in lines:
        if '=== ANSWER KEY ===' in line:
            answer_key_started = True
            continue
        
        if answer_key_started:
            answers_part.append(line)
        else:
            questions_part.append(line)
    
    # Handle ASSIGNMENT type
    if exam_type == 'assignment':
        task_num = 0
        for line in questions_part:
            line_stripped = line.strip()
            if line_stripped.startswith('Task '):
                task_num += 1
                tasks.append(line_stripped)
        
        return {
            'type': 'assignment',
            'tasks': tasks,
            'answers': '\n'.join(answers_part),
            'marks_config': marks_config
        }
    
    # Handle QUIZ and MID/FINAL types
    for line in questions_part:
        line_stripped = line.strip()
        
        if not line_stripped:
            continue
        
        # Detect section headers
        if line_stripped.startswith('SECTION'):
            if current_section:
                sections.append(current_section)
            
            section_name = line_stripped.split(':')[0].strip()
            current_section = {
                'name': section_name,
                'mcqs': [],
                'shorts': [],
                'longs': [],
                'scenarios': [],
                'total_marks': 0
            }
            current_qtype = None
        
        # Detect question type headers
        elif 'Multiple Choice' in line_stripped:
            current_qtype = 'mcqs'
        elif 'Short Answer' in line_stripped:
            current_qtype = 'shorts'
        elif 'Long Answer' in line_stripped:
            current_qtype = 'longs'
        elif 'Scenario' in line_stripped and current_section:
            current_qtype = 'scenarios'
        elif line_stripped and current_section and current_qtype:
            # Add question to current type
            if line_stripped[0].isdigit() or line_stripped.startswith('Task'):
                current_section[current_qtype].append(line_stripped)
    
    if current_section:
        sections.append(current_section)
    
    # Calculate total marks for each section
    for section in sections:
        section['total_marks'] = (
            len(section['mcqs']) * marks_config.get('mcq', 1) +
            len(section['shorts']) * marks_config.get('short', 5) +
            len(section['longs']) * marks_config.get('long', 10) +
            len(section['scenarios']) * marks_config.get('scenario', 10)
        )
    
    return {
        'type': 'exam',
        'sections': sections,
        'answers': '\n'.join(answers_part),
        'marks_config': marks_config
    }

# ------------------ FORMATTED DOCX EXPORT ------------------
# ------------------ FORMATTED DOCX EXPORT ------------------
def generate_docx(content: str, include_answers: bool, prompt: str = "", exam_type: str = "quiz"):
    """
    VERY SIMPLE EXPORT:
    - Ignore exam_type / marks / structure
    - Just write the raw exam.content line by line into a DOCX
    - This guarantees that NOTHING is dropped
    """
    doc = Document()

    # Optional: simple heading
    # doc.add_heading('TeachAssist Export', level=1)

    for line in content.splitlines():
        # Empty line => blank paragraph (keeps spacing)
        if line.strip() == "":
            doc.add_paragraph("")
        else:
            doc.add_paragraph(line)

    return doc

def generate_pdf(content: str, include_answers: bool, prompt: str = "", exam_type: str = "quiz"):
    """
    SUPER SIMPLE PDF EXPORT

    - No parsing / sections
    - Dumps the ENTIRE `content` string to PDF
    - Preserves line breaks, wraps long lines to page width
    """
    buffer = BytesIO()

    doc = SimpleDocTemplate(
        buffer,
        pagesize=A4,
        leftMargin=36,
        rightMargin=36,
        topMargin=36,
        bottomMargin=36,
    )

    styles = getSampleStyleSheet()
    body_style = ParagraphStyle(
        "Body",
        parent=styles["Normal"],
        fontName="Helvetica",
        fontSize=10,
        leading=13,
        wordWrap="LTR",   # wrap long lines
    )

    # Escape minimal HTML and convert newlines to <br/> for Paragraph
    safe_text = (
        content.replace("&", "&amp;")
               .replace("<", "&lt;")
               .replace(">", "&gt;")
    )
    safe_text = safe_text.replace("\r\n", "\n").replace("\r", "\n")
    safe_text = safe_text.replace("\n", "<br/>")

    story = [Paragraph(safe_text, body_style)]
    doc.build(story)

    buffer.seek(0)
    return buffer



# ------------------ SIMPLE RAW PDF EXPORT ------------------
@app.get("/api/download/{exam_id}")
def download_exam(
    exam_id: int,
    format: str,
    include_answers: bool = True,
    db: Session = Depends(get_db)
):
    exam = db.query(models.GeneratedExam).get(exam_id)
    if not exam:
        raise HTTPException(404, "Exam not found")

    exam_type = exam.exam_type if hasattr(exam, "exam_type") else "quiz"
    prompt = ""

    if format == "docx":
        doc = generate_docx(exam.content, include_answers, prompt, exam_type)
        buf = BytesIO()
        doc.save(buf)
        buf.seek(0)
        return StreamingResponse(
            buf,
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            headers={"Content-Disposition": "attachment; filename=exam.docx"}
        )

    if format == "pdf":
        pdf = generate_pdf(exam.content, include_answers, prompt, exam_type)
        return StreamingResponse(
            pdf,
            media_type="application/pdf",
            headers={"Content-Disposition": "attachment; filename=exam.pdf"}
        )

    raise HTTPException(400, "Invalid format")


# ------------------ LLaMA GENERATION ------------------
def generate_with_llama(prompt: str) -> str:
    """Generate with focused system prompt"""
    response = groq_client.chat.completions.create(
        model=GROQ_MODEL,
        messages=[
            {
                "role": "system",
                "content": "You are a precise university exam generator. Follow instructions exactly. Generate questions ONLY from provided context. Create realistic workplace scenarios when requested. Count questions carefully before outputting."
            },
            {
                "role": "user",
                "content": prompt
            }
        ],
        temperature=0.2,
        max_tokens=1800
    )
    return response.choices[0].message.content
# ===================== ASSIGNMENT HELPERS (add above the endpoint) =====================
ASSIGN_TOTAL_RE = re.compile(r"TOTAL_TASKS:\s*(\d+)", re.IGNORECASE)
ASSIGN_SCEN_RE  = re.compile(r"SCENARIO_TASKS:\s*(\d+)", re.IGNORECASE)

def extract_assignment_requirements(prompt: str) -> tuple[int, int]:
    t = ASSIGN_TOTAL_RE.search(prompt)
    s = ASSIGN_SCEN_RE.search(prompt)
    total = int(t.group(1)) if t else 0
    scenario = int(s.group(1)) if s else 0
    return total, scenario


def _parse_assignment_tasks(questions: str):
    """
    Split the questions-part into:
      - prefix_lines: everything before Task 1
      - tasks: list of {index, header, body}
    We accept headers like: 'Task 1:', 'Task 1 (5 marks):', etc.
    """
    lines = questions.splitlines()
    prefix_lines: list[str] = []
    tasks: list[dict] = []

    current_header: str | None = None
    current_index: int | None = None
    current_body_lines: list[str] = []

    for line in lines:
        m = re.match(r"\s*(Task\s+(\d+)[^:]*:)", line)
        if m:
            # starting a new Task block
            if current_header is None:
                # first task -> everything before is prefix
                prefix_lines = []
            else:
                # store previous task
                tasks.append({
                    "index": current_index,
                    "header": current_header,
                    "body": "\n".join(current_body_lines).rstrip()
                })

            current_header = m.group(1).rstrip()   # only up to the colon
            current_index = int(m.group(2))

            # anything after the colon on the same line becomes the first body line
            remainder = line[m.end():].rstrip()
            current_body_lines = [remainder] if remainder else []
        else:
            if current_header is None:
                prefix_lines.append(line)
            else:
                current_body_lines.append(line)

    # flush last task
    if current_header is not None:
        tasks.append({
            "index": current_index,
            "header": current_header,
            "body": "\n".join(current_body_lines).rstrip()
        })

    return prefix_lines, tasks


def _normalize_to_scenario(body: str) -> str:
    """
    Convert the body of a task into:
        Scenario: ...
        Task: ...
    If it's already in proper Scenario+Task format, keep as is.
    """
    has_scenario = re.search(r"(?im)^\s*Scenario\s*:", body) is not None
    has_task     = re.search(r"(?im)^\s*Task\s*:", body)     is not None

    # already proper scenario → just tidy
    if has_scenario and has_task:
        return body.strip() + "\n"

    lines = [ln for ln in body.strip().splitlines() if ln.strip()]
    first_line = lines[0] if lines else ""
    rest_text  = "\n".join(lines[1:]).strip()

    scenario_stub = (
        "Scenario: Consider a realistic workplace/project situation "
        "relevant to the uploaded lecture content."
    )
    task_line = (
        f"Task: {first_line.strip()}" if first_line
        else "Task: Respond to the scenario."
    )

    if rest_text:
        task_line += f"\n{rest_text}"

    return f"{scenario_stub}\n{task_line}\n"


def force_assignment_scenarios(text: str, total_tasks: int, scenario_tasks: int) -> str:
    """
    Ensure the last `scenario_tasks` tasks are SCENARIO style.
    - If SCENARIO_TASKS == TOTAL_TASKS → ALL tasks become scenarios.
    - Only modifies the questions part (before '=== ANSWER KEY ===').
    """
    if total_tasks <= 0 or scenario_tasks <= 0:
        return text

    parts = text.split("=== ANSWER KEY ===", 1)
    questions = parts[0]
    answers   = parts[1] if len(parts) > 1 else ""

    prefix_lines, tasks = _parse_assignment_tasks(questions)
    if not tasks:
        # nothing recognized as tasks, bail out
        return text

    present_indices = [t["index"] for t in tasks]

    # Clamp scenario_tasks to actually existing tasks
    scenario_tasks = min(scenario_tasks, len(present_indices))

    # We always convert the LAST N tasks into scenarios
    target_indices = present_indices[-scenario_tasks:]
    target_set = set(target_indices)

    # Rebuild the questions part
    out_lines: list[str] = []
    out_lines.extend(prefix_lines)

    for t in tasks:
        body = t["body"]
        if t["index"] in target_set:
            body = _normalize_to_scenario(body)

        block = t["header"] + "\n"
        if body:
            block += body
        out_lines.append(block)

    normalized_questions = "\n".join(out_lines).strip()

    if answers:
        return f"{normalized_questions}\n\n=== ANSWER KEY ===\n{answers.lstrip()}"
    return normalized_questions

def trim_assignment_tasks(text: str, total_tasks: int) -> str:
    """
    Keep only Task 1..TOTAL_TASKS in both QUESTIONS and ANSWER KEY.
    Drop any Task N where N > TOTAL_TASKS.
    """
    if not total_tasks or total_tasks <= 0:
        return text

    parts = text.split("=== ANSWER KEY ===", 1)
    questions = parts[0]
    answers   = parts[1] if len(parts) > 1 else ""

    # ---- trim questions part using _parse_assignment_tasks ----
    prefix_lines, tasks = _parse_assignment_tasks(questions)
    if not tasks or len(tasks) <= total_tasks:
        # nothing to trim
        return text

    kept_tasks = tasks[:total_tasks]

    out_q_lines: list[str] = []
    out_q_lines.extend(prefix_lines)

    for t in kept_tasks:
        block = t["header"]
        body  = t["body"]
        out_q_lines.append(block)
        if body:
            out_q_lines.append(body)
        # blank line between tasks (optional but keeps layout nice)
        out_q_lines.append("")

    normalized_questions = "\n".join(out_q_lines).strip()

    # ---- trim answer key (if present) ----
    if not answers:
        return normalized_questions

    filtered_answer_lines: list[str] = []
    for line in answers.splitlines():
        m = re.match(r"\s*Task\s+(\d+)\b", line)
        if m:
            idx = int(m.group(1))
            if idx > total_tasks:
                # skip answers for extra tasks
                continue
        filtered_answer_lines.append(line)

    normalized_answers = "\n".join(filtered_answer_lines).strip()

    return f"{normalized_questions}\n\n=== ANSWER KEY ===\n{normalized_answers}\n"

# =========================================================
# GENERATE EXAM (MAIN ENDPOINT)
# =========================================================
@app.post("/api/generate-exam", response_model=schemas.ExamOut)
async def generate_exam(
    exam_type: str = Form(...),
    prompt: str = Form(...),
    teacher_prompt: str = Form(""),
    files: List[UploadFile] = File(...),
    db: Session = Depends(get_db)
):
    try:
        # normalize exam type once
        exam_type_lower = exam_type.lower()

        # ------------------ Ensure teacher ------------------
        teacher = db.query(models.Teacher).first()
        if not teacher:
            teacher = models.Teacher(
                teacher_uid="TCH-DEFAULT",
                username="default.teacher",
                teacher_name="default_teacher",
                password=hash_password("1234")
            )
            db.add(teacher)
            db.commit()
            db.refresh(teacher)

        # ------------------ Extract text from files ------------------
        all_documents = []
        for idx, file in enumerate(files):
            text = extract_text_from_file(file)
            if text.strip():
                all_documents.append({
                    "doc_id": idx + 1,
                    "content": text
                })

        if not all_documents:
            raise HTTPException(400, "Empty lecture content")

        # ------------------ Chunk + Store in RAG ------------------
        all_chunks = []
        all_metadata = []

        for doc in all_documents:
            chunks = chunk_text(doc["content"])
            all_chunks.extend(chunks)
            all_metadata.extend(
                [{"teacher_id": teacher.id, "doc_id": doc["doc_id"]}] * len(chunks)
            )

        rag.add_documents(texts=all_chunks, metadatas=all_metadata)

        # ------------------ Retrieve balanced context from RAG ------------------
        doc_chunks_map = defaultdict(list)
        docs_count = len(all_documents)
        TOTAL_K = 6
        k_per_doc = max(1, TOTAL_K // docs_count)

        for doc in all_documents:
            chunks = rag.search(
                query="Generate exam questions strictly from this document",
                top_k=k_per_doc,
                filter={"doc_id": doc["doc_id"]}
            )
            doc_chunks_map[doc["doc_id"]].extend(chunks)

        if not doc_chunks_map:
            raise HTTPException(400, "No relevant content retrieved from RAG")

        # Build context with document separation
        context_parts = []
        for doc_id in sorted(doc_chunks_map.keys()):
            chunks = doc_chunks_map[doc_id]
            joined_chunks = clip_text("\n\n".join(chunks), 1200)
            context_parts.append(
                f"""===== DOCUMENT {doc_id} =====
{joined_chunks}"""
            )

        context = clip_text("\n\n".join(context_parts), 5000)
        compact_prompt = clip_text(prompt.strip(), 4000)
        compact_teacher_prompt = clip_text(teacher_prompt.strip(), 1000)

   # ------------------ Build final prompt (FRONTEND DRIVEN) ------------------

        final_prompt = f"""
You are an AI exam generator for teachers.

You MUST strictly follow the instructions provided in the UI PROMPT.
You MUST generate questions ONLY from the lecture context below.

================ LECTURE CONTEXT =================
{context}

================ UI PROMPT (STRICT RULES) =================
{compact_prompt}

================ TEACHER NOTES =================
{compact_teacher_prompt}
"""

        # Single LLaMA call for ALL exam types
        raw_exam_text = generate_with_llama(final_prompt)
        exam_text = clean_output(raw_exam_text)

                # Extra post-processing ONLY for assignments
        if exam_type_lower == "assignment":
            total_tasks, scenario_tasks = extract_assignment_requirements(prompt)

            # 1) Enforce the number of tasks (drop Task 6,7,8,9, …)
            if total_tasks > 0:
                exam_text = trim_assignment_tasks(exam_text, total_tasks)

            # 2) Force the last N tasks to be scenarios (if any)
            if scenario_tasks > 0:
                exam_text = force_assignment_scenarios(
                    exam_text,
                    total_tasks,
                    scenario_tasks
                )

        # ------------------ Save to database ------------------
        exam = models.GeneratedExam(
            teacher_id=teacher.id,
            exam_type=exam_type,   # original string ('quiz', 'assignment', 'midterm')
            content=exam_text
        )
        db.add(exam)
        db.commit()
        db.refresh(exam)

        return exam

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

# ------------------ DOWNLOAD ENDPOINT ------------------
@app.get("/api/download/{exam_id}")
def download_exam(
    exam_id: int,
    format: str,
    include_answers: bool = True,
    db: Session = Depends(get_db)
):
    exam = db.query(models.GeneratedExam).get(exam_id)
    if not exam:
        raise HTTPException(404, "Exam not found")

    exam_type = exam.exam_type if hasattr(exam, "exam_type") else "quiz"
    prompt = ""

    if format == "docx":
        doc = generate_docx(exam.content, include_answers, prompt, exam_type)
        buf = BytesIO()
        doc.save(buf)
        buf.seek(0)
        return StreamingResponse(
            buf,
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            headers={"Content-Disposition": "attachment; filename=exam.docx"}
        )

    if format == "pdf":
        pdf = generate_pdf(exam.content, include_answers, prompt, exam_type)
        return StreamingResponse(
            pdf,
            media_type="application/pdf",
            headers={"Content-Disposition": "attachment; filename=exam.pdf"}
        )

    raise HTTPException(400, "Invalid format")
